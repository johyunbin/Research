"""
pptx_helpers.py — Samsung Research Portfolio PPT design system module.

Samsung brand colors, slide dimensions, font constants, and reusable layout functions
for building the 18-slide portfolio presentation.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Samsung Brand Color Constants
# ---------------------------------------------------------------------------
NAVY        = RGBColor(0x13, 0x28, 0x9F)
ACCENT_BLUE = RGBColor(0x3D, 0x7D, 0xDE)
SKY_BLUE    = RGBColor(0x06, 0x89, 0xD8)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY        = RGBColor(0x5B, 0x71, 0x8D)
GRAY2       = RGBColor(0x75, 0x78, 0x7B)
LIGHT_BG    = RGBColor(0xF0, 0xF4, 0xF8)
WARM_GRAY   = RGBColor(0xE7, 0xE6, 0xE2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)

# ---------------------------------------------------------------------------
# Slide Dimensions  (16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Margins
# ---------------------------------------------------------------------------
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.5)

# ---------------------------------------------------------------------------
# Font Constants
# ---------------------------------------------------------------------------
FONT_TITLE = "맑은 고딕"
FONT_BODY  = "맑은 고딕"
FONT_EN    = "Segoe UI"


# ---------------------------------------------------------------------------
# Presentation Factory
# ---------------------------------------------------------------------------

def new_presentation() -> Presentation:
    """Return a blank 16:9 Presentation with Samsung dimensions."""
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


# ---------------------------------------------------------------------------
# Slide Helpers
# ---------------------------------------------------------------------------

def add_blank_slide(prs: Presentation):
    """Add and return a fully blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]  # index 6 is the blank layout
    return prs.slides.add_slide(blank_layout)


def set_slide_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Shape Primitives
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height,
             fill_color: RGBColor = None,
             line_color: RGBColor = None,
             line_width_pt: float = 0.0):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = shape.fill
    if fill_color is not None:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color is not None:
        line.color.rgb = line_color
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()

    return shape


def add_accent_bar(slide, left, top, width, height=Inches(0.04),
                   color: RGBColor = None):
    """Add a thin horizontal accent bar (defaults to NAVY)."""
    color = color or NAVY
    return add_rect(slide, left, top, width, height, fill_color=color)


# ---------------------------------------------------------------------------
# Text Helpers
# ---------------------------------------------------------------------------

def add_textbox(slide, left, top, width, height, text: str,
                font_size: float = 12,
                font_color: RGBColor = None,
                bold: bool = False,
                alignment=PP_ALIGN.LEFT,
                font_name: str = None,
                anchor=MSO_ANCHOR.TOP) -> object:
    """
    Add a single-paragraph textbox and return the shape.

    Parameters
    ----------
    slide       : slide object
    left/top/width/height : position & size (Emu / Inches values)
    text        : string content
    font_size   : in points
    font_color  : RGBColor (defaults to DARK_TEXT)
    bold        : bool
    alignment   : PP_ALIGN constant
    font_name   : override font family (defaults to FONT_BODY)
    anchor      : MSO_ANCHOR vertical anchor
    """
    font_color = font_color if font_color is not None else DARK_TEXT
    font_name  = font_name  or FONT_BODY

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = alignment

    run = p.add_run()
    run.text = text

    fnt = run.font
    fnt.size  = Pt(font_size)
    fnt.bold  = bold
    fnt.color.rgb = font_color
    fnt.name  = font_name

    return txBox


def add_multiline_textbox(slide, left, top, width, height,
                          lines,
                          font_size: float = 11,
                          font_color: RGBColor = None,
                          line_spacing: float = 1.15,
                          font_name: str = None,
                          alignment=PP_ALIGN.LEFT) -> object:
    """
    Add a textbox with multiple lines/runs.

    Parameters
    ----------
    lines : list of items, each either:
            - str  → rendered with defaults
            - tuple (text, bold, color, size)  → all four fields
            - tuple (text, bold, color)  → size falls back to font_size
            - tuple (text, bold)         → color & size fall back
    line_spacing : paragraph line spacing multiplier
    """
    font_color = font_color if font_color is not None else DARK_TEXT
    font_name  = font_name  or FONT_BODY

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    first = True
    for item in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        para.alignment = alignment

        # Parse line spec
        if isinstance(item, str):
            t, b, c, s = item, False, font_color, font_size
        elif isinstance(item, (list, tuple)):
            if len(item) == 4:
                t, b, c, s = item
            elif len(item) == 3:
                t, b, c = item; s = font_size
            elif len(item) == 2:
                t, b = item; c = font_color; s = font_size
            else:
                t = item[0]; b = False; c = font_color; s = font_size
        else:
            t, b, c, s = str(item), False, font_color, font_size

        # Line spacing
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        pPr = para._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, _qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))

        run = para.add_run()
        run.text = t

        fnt = run.font
        fnt.size  = Pt(s)
        fnt.bold  = b
        fnt.color.rgb = c if c is not None else font_color
        fnt.name  = font_name

    return txBox


# ---------------------------------------------------------------------------
# Composite Components
# ---------------------------------------------------------------------------

def add_section_title(slide, left, top, width,
                      title_text: str,
                      subtitle_text: str = "",
                      title_size: float = 28,
                      subtitle_size: float = 14,
                      bar_height=None,
                      bar_color: RGBColor = None):
    """
    Add a section title block: accent bar + big title + optional subtitle.

    Returns (bar_shape, title_shape, subtitle_shape or None).
    """
    bar_height = bar_height or Inches(0.055)
    bar_color  = bar_color  or NAVY

    bar = add_accent_bar(slide, left, top, width, bar_height, color=bar_color)

    title_top = top + bar_height + Inches(0.08)
    title_h   = Inches(0.55)
    t_shape   = add_textbox(
        slide, left, title_top, width, title_h,
        title_text,
        font_size=title_size,
        font_color=NAVY,
        bold=True,
        font_name=FONT_BODY,
    )

    s_shape = None
    if subtitle_text:
        sub_top = title_top + title_h
        s_shape = add_textbox(
            slide, left, sub_top, width, Inches(0.35),
            subtitle_text,
            font_size=subtitle_size,
            font_color=GRAY,
            bold=False,
            font_name=FONT_BODY,
        )

    return bar, t_shape, s_shape


def add_metric_card(slide, left, top, width, height,
                    number_text: str,
                    label_text: str,
                    bg_color: RGBColor = None,
                    number_size: float = 32,
                    label_size: float = 11):
    """
    A light-BG card with a big Navy number on top and a small gray label below.

    Returns the background rect shape.
    """
    bg_color = bg_color or LIGHT_BG

    bg = add_rect(slide, left, top, width, height, fill_color=bg_color)

    pad     = Inches(0.15)
    inner_w = width - pad * 2

    # Number
    num_h   = height * 0.55
    add_textbox(
        slide,
        left + pad,
        top  + pad,
        inner_w,
        num_h,
        number_text,
        font_size=number_size,
        font_color=NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name=FONT_EN,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Label
    lbl_top = top + num_h
    lbl_h   = height - num_h - pad
    add_textbox(
        slide,
        left + pad,
        lbl_top,
        inner_w,
        lbl_h,
        label_text,
        font_size=label_size,
        font_color=GRAY2,
        bold=False,
        alignment=PP_ALIGN.CENTER,
        font_name=FONT_BODY,
        anchor=MSO_ANCHOR.TOP,
    )

    return bg


def add_implication_box(slide, left, top, width, height,
                        lines,
                        bg_color: RGBColor = None,
                        font_size: float = 11,
                        line_spacing: float = 1.3):
    """
    Navy background box with white text lines.

    lines : list of str or tuple specs accepted by add_multiline_textbox.
    Returns the background rect shape.
    """
    bg_color = bg_color or NAVY

    bg = add_rect(slide, left, top, width, height, fill_color=bg_color)

    pad    = Inches(0.2)
    inner_w = width  - pad * 2
    inner_h = height - pad * 2

    # Normalise lines so all colors default to WHITE
    normalised = []
    for item in lines:
        if isinstance(item, str):
            normalised.append((item, False, WHITE, font_size))
        elif isinstance(item, (list, tuple)):
            if len(item) == 4:
                normalised.append(item)
            elif len(item) == 3:
                normalised.append((item[0], item[1], item[2], font_size))
            elif len(item) == 2:
                normalised.append((item[0], item[1], WHITE, font_size))
            else:
                normalised.append((str(item[0]), False, WHITE, font_size))
        else:
            normalised.append((str(item), False, WHITE, font_size))

    add_multiline_textbox(
        slide,
        left  + pad,
        top   + pad,
        inner_w,
        inner_h,
        normalised,
        font_size=font_size,
        font_color=WHITE,
        line_spacing=line_spacing,
        font_name=FONT_BODY,
        alignment=PP_ALIGN.LEFT,
    )

    return bg


def add_image_safe(slide, image_path: str, left, top, width, height,
                   placeholder_color: RGBColor = None,
                   placeholder_text: str = ""):
    """
    Add an image to the slide.  If the file is missing, add a colored
    placeholder rectangle with optional label text instead.

    Returns the added shape.
    """
    if image_path and os.path.isfile(image_path):
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        return pic
    else:
        ph_color = placeholder_color or WARM_GRAY
        rect = add_rect(slide, left, top, width, height,
                        fill_color=ph_color,
                        line_color=GRAY2,
                        line_width_pt=0.75)
        label = placeholder_text or (os.path.basename(image_path) if image_path else "Image")
        add_textbox(
            slide, left, top, width, height,
            label,
            font_size=10,
            font_color=GRAY,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        return rect
