#!/usr/bin/env python3
"""
Hanyang University Faculty Interview Presentation
"Intelligent Sound Environment Systems for Human-Centered Built Environments"
Hyun In Jo, Ph.D.

20 slides, 15-minute presentation
Design: Clean, professional, navy+white+accent blue
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Design Constants ───
NAVY = RGBColor(0x0A, 0x1E, 0x3D)       # Dark navy background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)  # Bright accent
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
GOLD = RGBColor(0xD4, 0xA5, 0x37)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ─── Helper Functions ───
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * (line_spacing - 1))
    return tf


def add_paragraph(tf, text, font_size=16, color=DARK_TEXT, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri",
                  space_before=0, space_after=4, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    return p


def add_accent_bar(slide, top, width=Inches(1.5), color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def make_section_header(slide, section_num, title, subtitle=""):
    """Navy background section divider slide"""
    add_bg(slide, NAVY)
    # Section number
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(2), Inches(1),
                 f"PART {section_num}", font_size=20, color=ACCENT_TEAL,
                 bold=True, font_name="Calibri")
    # Title
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(2),
                 title, font_size=40, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(1),
                     subtitle, font_size=18, color=ACCENT_TEAL)


def make_content_slide(slide, title, show_line=True):
    """White background content slide with title bar"""
    add_bg(slide, WHITE)
    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0), Inches(0), SLIDE_W, Pt(6))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_BLUE
    top_bar.line.fill.background()
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 title, font_size=28, color=NAVY, bold=True)
    if show_line:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.8), Inches(1.05), Inches(11.7), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()


def add_slide_number(slide, num, total=20):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# S1. TITLE SLIDE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Decorative accent line
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(1.8), Inches(2), Pt(4))
deco.fill.solid()
deco.fill.fore_color.rgb = ACCENT_TEAL
deco.line.fill.background()

# Title
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
             "Intelligent Sound Environment Systems\nfor Human-Centered Built Environments",
             font_size=36, color=WHITE, bold=True, line_spacing=1.3)

# Subtitle - position info
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(8), Inches(0.6),
             "Faculty Position in Intelligent Building Environment Systems",
             font_size=18, color=ACCENT_TEAL, bold=False)

# Separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(4.7), Inches(5), Pt(1))
sep.fill.solid()
sep.fill.fore_color.rgb = RGBColor(0x33, 0x55, 0x77)
sep.line.fill.background()

# Name
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.6),
             "Hyun In Jo, Ph.D.", font_size=28, color=WHITE, bold=True)

# Credentials
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(8), Inches(0.4),
             "LEED AP BD+C  |  LEED AP ID+C  |  WELL AP  |  ADsP",
             font_size=14, color=ACCENT_TEAL)

# Affiliation
add_text_box(slide, Inches(0.8), Inches(6.1), Inches(10), Inches(0.8),
             "Department of Architectural Engineering, Hanyang University",
             font_size=16, color=RGBColor(0x99, 0xAA, 0xBB))


# ═══════════════════════════════════════════════════
# S2. TABLE OF CONTENTS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Table of Contents")
add_slide_number(slide, 2)

sections = [
    ("01", "Introduction", "The Evolution  |  Brief Introduction"),
    ("02", "Research", "Research Philosophy  |  Key Research (Buildings · Cities · Products · Health)\nGlobal Network  |  Industry Experience"),
    ("03", "Future Research Plan", "SENSE Lab Vision  |  Phase 1-2-3  |  Funding Roadmap"),
    ("04", "Education", "Teaching Philosophy  |  Teaching Plan  |  Student Mentoring"),
    ("05", "Contribution", "Department  |  Industry  |  Global"),
]

y_start = Inches(1.5)
for i, (num, title, desc) in enumerate(sections):
    y = y_start + Inches(i * 1.1)
    # Number box
    num_box = add_shape(slide, Inches(0.8), y, Inches(0.8), Inches(0.7), fill_color=NAVY)
    tf_num = num_box.text_frame
    tf_num.paragraphs[0].text = num
    tf_num.paragraphs[0].font.size = Pt(20)
    tf_num.paragraphs[0].font.color.rgb = WHITE
    tf_num.paragraphs[0].font.bold = True
    tf_num.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Section title
    add_text_box(slide, Inches(1.8), y, Inches(3), Inches(0.7),
                 title, font_size=22, color=NAVY, bold=True)

    # Description
    add_text_box(slide, Inches(4.8), y + Inches(0.05), Inches(7.5), Inches(0.7),
                 desc, font_size=14, color=MID_GRAY)


# ═══════════════════════════════════════════════════
# S3. THE EVOLUTION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "The Evolution of Building Sound Environment Research")
add_slide_number(slide, 3)

# Three generation boxes
gen_data = [
    ("Generation 1", "Noise Control", "Physical metrics\ndB, RT60, NC curves\n\n\"How loud is it?\"",
     RGBColor(0xE8, 0xE8, 0xE8), DARK_TEXT),
    ("Generation 2", "Soundscape", "Perceptual experience\nISO 12913, VR, Bio-signals\n\n\"How do people feel?\"",
     RGBColor(0xD6, 0xEA, 0xF8), NAVY),
    ("Generation 3", "Intelligent\nSound Environment", "Data-driven optimization\nAI, IoT, Digital Twin\n\n\"How to optimize\nautomatically?\"",
     NAVY, WHITE),
]

for i, (gen, title, desc, bg_col, txt_col) in enumerate(gen_data):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.5)

    # Box
    box = add_shape(slide, x, y, Inches(3.8), Inches(4.5), fill_color=bg_col)
    box.shadow.inherit = False

    # Gen label
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.4),
                 gen, font_size=13, color=ACCENT_BLUE if i < 2 else ACCENT_TEAL, bold=True)

    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(3.2), Inches(1.2),
                 title, font_size=24, color=txt_col, bold=True, line_spacing=1.1)

    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.2), Inches(2.2),
                 desc, font_size=14, color=txt_col if i == 2 else MID_GRAY, line_spacing=1.3)

    # Arrow between boxes
    if i < 2:
        arrow_x = x + Inches(3.85)
        add_text_box(slide, arrow_x, Inches(3.2), Inches(0.3), Inches(0.5),
                     "→", font_size=24, color=ACCENT_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)

# Bottom statement
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
             "\"I have researched across all three generations — and I bring the capability to lead Generation 3.\"",
             font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# S4. BRIEF INTRODUCTION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Brief Introduction")
add_slide_number(slide, 4)

# Journey Timeline
timeline_items = [
    ("2013 – 2016", "B.S. Architectural Engineering", "Hanyang University (Summa Cum Laude, Early Graduation)", ACCENT_BLUE),
    ("2016 – 2022", "Ph.D. Architectural Environmental Engineering", "Hanyang University (GPA 4.39/4.5)", ACCENT_BLUE),
    ("2018", "Visiting Researcher", "Sorbonne University, Paris, France", ACCENT_TEAL),
    ("2022", "Post-doctoral Researcher (NST Fellow)", "Korea Institute of Civil Eng. & Building Tech. (KICT)", ACCENT_TEAL),
    ("2022 – Present", "Senior Research Engineer", "Hyundai Motor Company, NVH Test Research Lab", GOLD),
]

for i, (period, role, org, dot_color) in enumerate(timeline_items):
    y = Inches(1.4) + Inches(i * 0.75)
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.08), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    # Vertical line
    if i < len(timeline_items) - 1:
        vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(1.07), y + Inches(0.26), Pt(2), Inches(0.55))
        vline.fill.solid()
        vline.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        vline.line.fill.background()
    # Period
    add_text_box(slide, Inches(1.4), y - Inches(0.05), Inches(2), Inches(0.35),
                 period, font_size=13, color=MID_GRAY, bold=True)
    # Role
    add_text_box(slide, Inches(3.4), y - Inches(0.05), Inches(4.5), Inches(0.35),
                 role, font_size=14, color=NAVY, bold=True)
    # Org
    add_text_box(slide, Inches(3.4), y + Inches(0.22), Inches(5), Inches(0.35),
                 org, font_size=12, color=MID_GRAY)

# Right side - Stats Dashboard
stats_x = Inches(8.5)
stats = [
    ("26", "SCI(E) Papers\n(21 First Author)"),
    ("18", "h-index"),
    ("6", "Patents\n(incl. 1 US)"),
    ("12", "Awards"),
    ("$5M+", "Funded\nResearch"),
    ("4,565", "Concert Hall\nSeats Designed"),
]

for i, (num, label) in enumerate(stats):
    row = i // 3
    col = i % 3
    x = stats_x + Inches(col * 1.5)
    y = Inches(1.5) + Inches(row * 2.2)

    add_text_box(slide, x, y, Inches(1.4), Inches(0.7),
                 num, font_size=28, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.65), Inches(1.4), Inches(0.8),
                 label, font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.2)


# ═══════════════════════════════════════════════════
# S5. RESEARCH PHILOSOPHY
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Research Philosophy")
add_slide_number(slide, 5)

# Vision statement
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.8),
             "\"Designing human-centric acoustic environments by integrating\nperception science, architecture, and intelligent technology\"",
             font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER,
             line_spacing=1.3)

# Three pillars
pillars = [
    ("Perception over Physics",
     "Evaluate environments through\nhuman experience, not just dB.\n\nPsychoacoustics, bio-signals\n(EEG, HRV, eye-tracking),\nand subjective assessment\nreplace purely physical metrics.",
     "🎯"),
    ("Interdisciplinary Integration",
     "Bridge architectural acoustics\nwith cognitive psychology,\nAI/data science, and urban design.\n\nOne research methodology\napplied across buildings, cities,\nproducts, and healthcare.",
     "🔗"),
    ("Lab to Life",
     "Research must reach real buildings,\ncities, and products.\n\n3 concert halls designed,\nEV sound in mass production,\ndigital therapeutics validated\nin clinical settings.",
     "🏗️"),
]

for i, (title, desc, icon) in enumerate(pillars):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.5)

    # Pillar box
    box = add_shape(slide, x, y, Inches(3.8), Inches(4.5), fill_color=LIGHT_BLUE_BG)

    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.6),
                 title, font_size=18, color=NAVY, bold=True)

    # Accent bar under title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Inches(0.3), y + Inches(0.85), Inches(1.5), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.2), Inches(3.0),
                 desc, font_size=13, color=MID_GRAY, line_spacing=1.35)


# ═══════════════════════════════════════════════════
# S6. RESEARCH OVERVIEW
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Research Overview — One Method, Four Impacts")
add_slide_number(slide, 6)

# Pipeline: SENSE → UNDERSTAND → PREDICT → APPLY
pipeline_steps = [
    ("SENSE", "VR/AR\nAuralization"),
    ("UNDERSTAND", "Psychoacoustic\nEvaluation"),
    ("PREDICT", "AI / ML\nModeling"),
    ("APPLY", "Design &\nProduct"),
]

for i, (step, desc) in enumerate(pipeline_steps):
    x = Inches(0.6) + Inches(i * 3.2)
    y = Inches(1.4)
    # Box
    box = add_shape(slide, x, y, Inches(2.6), Inches(1.5), fill_color=NAVY)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = ACCENT_TEAL
    p2.alignment = PP_ALIGN.CENTER

    if i < 3:
        add_text_box(slide, x + Inches(2.65), y + Inches(0.35), Inches(0.5), Inches(0.5),
                     "→", font_size=22, color=ACCENT_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)

# Four Impact Areas
impacts = [
    ("BUILDINGS", "Concert halls, apartments,\noffice environments", "12 SCI papers"),
    ("CITIES", "Urban parks, streets,\nsoundscape design", "7 SCI papers"),
    ("PRODUCTS", "EV AVAS, vehicle NVH,\nmass production", "2 SCI + JASA (review)"),
    ("HEALTH", "Pneumonia diagnosis,\ndigital therapeutics", "3 SCI + 1 TT"),
]

for i, (area, desc, output) in enumerate(impacts):
    x = Inches(0.6) + Inches(i * 3.2)
    y = Inches(3.5)

    box = add_shape(slide, x, y, Inches(2.6), Inches(3.3),
                    fill_color=LIGHT_BLUE_BG)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.2), Inches(0.5),
                 area, font_size=16, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Inches(0.5), y + Inches(0.7), Inches(1.6), Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.3)

    add_text_box(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.2), Inches(0.5),
                 output, font_size=11, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# S7. IMPACT 1: BUILDINGS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Impact 1: Buildings — Comfortable Built Environments")
add_slide_number(slide, 7)

# Left column - Concert Hall Design
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "Architectural Acoustics Design", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(1.75))

items_left = [
    "• 3 Concert Halls Designed: IFEZ Art Center (1,800 seats),\n  Bucheon Culture Arts Center (1,450), Pyeongtaek (1,315)",
    "• Sound diffusion optimization using scale models + simulations",
    "  → Jo & Jeon (2022) J. Building Engineering, Q1, IF 6.4",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(2.0),
                  items_left[0], font_size=13, color=DARK_TEXT, line_spacing=1.4)
for item in items_left[1:]:
    add_paragraph(tf, item, font_size=13, color=DARK_TEXT if not item.startswith("  →") else ACCENT_BLUE,
                  bold=item.startswith("  →"), space_after=6)

# Left column - Indoor Environment
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.5),
             "Indoor Sound Environment Evaluation", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(4.25))

items_indoor = [
    "• Floor Impact Noise: VR-based evaluation criteria",
    "  → Criterion lowered by 6-7 dB vs. existing standards",
    "  → Jo & Jeon (2019) Building & Environment, Q1, Cited 50",
    "",
    "• Open-Plan Office: Comfort-Content balance model",
    "  → Trade-off between preference and productivity discovered",
    "  → 2 papers in Building & Environment, Q1 (IF 7.4)",
    "",
    "• Water Supply/Drainage Noise: VR annoyance assessment",
    "  → Jeon, Jo* et al. (2019) Applied Acoustics, Q1",
]
tf2 = add_text_box(slide, Inches(0.8), Inches(4.45), Inches(5.5), Inches(2.8),
                   items_indoor[0], font_size=13, color=DARK_TEXT, line_spacing=1.3)
for item in items_indoor[1:]:
    c = ACCENT_BLUE if item.startswith("  →") else DARK_TEXT
    b = item.startswith("  →")
    add_paragraph(tf2, item, font_size=13 if not b else 12, color=c, bold=False, space_after=2)

# Right side - Key message
msg_box = add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
                    fill_color=LIGHT_BLUE_BG)
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.0), Inches(0.5),
             "Key Contribution", font_size=16, color=NAVY, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5.0), Inches(4.0),
             "Established human perception-based\nevaluation criteria for building acoustics,\nreplacing conventional physical metrics.\n\n"
             "VR technology enables reproducible\nlaboratory experiments that closely match\nreal-world acoustic experiences.\n\n"
             "Results directly applied to\nKorean building noise standards\nand architectural design practice.",
             font_size=14, color=MID_GRAY, line_spacing=1.4)


# ═══════════════════════════════════════════════════
# S8. IMPACT 2: CITIES
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Impact 2: Cities — Pleasant Urban Sound Environments")
add_slide_number(slide, 8)

# Soundscape Assessment
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "Soundscape Assessment & Design", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(1.75))

city_items = [
    "• ISO 12913 Framework: Quantitative + qualitative protocols",
    "  → Jo & Jeon (2021) Sustainable Cities & Society, Q1, IF 11.7",
    "",
    "• Audio-Visual Interaction: Visual 90% contribution,",
    "  yet sound environment determines overall satisfaction",
    "  → Jeon & Jo* (2020) Building & Environment, Cited 240",
    "",
    "• Soundscape Design Index: Revised GSI/RSI/GLI/RLI",
    "  SEM model: \"Secure pleasantness, control eventfulness\"",
    "  → Jo & Jeon (2021) Building & Environment, Cited 103",
    "",
    "• Urban Behavior: Human presence → dynamic soundscape",
    "  → Jo & Jeon (2020) Landscape & Urban Planning, Cited 111",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(6.0), Inches(4.0),
                  city_items[0], font_size=13, color=DARK_TEXT, line_spacing=1.25)
for item in city_items[1:]:
    c = ACCENT_BLUE if item.startswith("  →") else DARK_TEXT
    add_paragraph(tf, item, font_size=12 if item.startswith("  →") else 13,
                  color=c, space_after=2)

# Right - International Collaboration
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5),
             "International Collaboration", font_size=18, color=NAVY, bold=True)

collabs = [
    ("UCL, London", "Prof. Jian Kang", "Soundscape descriptors\n18-language validation (SATP)"),
    ("Sorbonne, Paris", "Prof. J-D. Polack", "Paris-Seoul urban soundscape\ncomparative study (STAR)"),
    ("RWTH Aachen", "Prof. M. Vorlaender", "VR spatial audio rendering\n& ecological validity"),
]

for i, (inst, prof, desc) in enumerate(collabs):
    y = Inches(2.0) + Inches(i * 1.7)
    box = add_shape(slide, Inches(7.2), y, Inches(5.3), Inches(1.4), fill_color=LIGHT_BLUE_BG)
    add_text_box(slide, Inches(7.5), y + Inches(0.1), Inches(2.5), Inches(0.4),
                 inst, font_size=14, color=NAVY, bold=True)
    add_text_box(slide, Inches(10.0), y + Inches(0.1), Inches(2.3), Inches(0.4),
                 prof, font_size=12, color=MID_GRAY)
    add_text_box(slide, Inches(7.5), y + Inches(0.5), Inches(4.8), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════
# S9. IMPACT 3: PRODUCTS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Impact 3: Products — From Research to Mass Production")
add_slide_number(slide, 9)

# AVAS section
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "EV AVAS Sound Design (Hyundai Motor Company)", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(1.75))

avas_items = [
    "• AVAS Brand Sound 2.0 Design & Tuning",
    "  Hyundai / Kia / Genesis (EV3, IONIQ 5)",
    "",
    "• Soundscape-based competitor AVAS evaluation",
    "  VR auralization with ambisonic reproduction",
    "",
    "• Sound environment quantification methodology (TDP)",
    "  for electric vehicle cabins using architectural acoustic theory",
    "",
    "• Company-wide AVAS expert committee member",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(2.8),
                  avas_items[0], font_size=13, color=DARK_TEXT, line_spacing=1.3)
for item in avas_items[1:]:
    add_paragraph(tf, item, font_size=13, color=DARK_TEXT, space_after=2)

# JASA Paper
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5),
             "AVAS × Urban Soundscape (JASA, under review)", font_size=16, color=ACCENT_BLUE, bold=True)
jasa_items = [
    "• 43 AVAS sounds from 17 EV models recorded & evaluated",
    "• ISO 12913 Pleasantness-Eventfulness assessment",
    "• Key finding: AVAS dissimilar to ICE engine sounds",
    "  → significantly higher pleasantness ratings",
    "• \"AVAS as a design tool for urban soundscape quality\"",
]
tf2 = add_text_box(slide, Inches(0.8), Inches(4.95), Inches(5.5), Inches(2.0),
                   jasa_items[0], font_size=12, color=DARK_TEXT, line_spacing=1.3)
for item in jasa_items[1:]:
    add_paragraph(tf2, item, font_size=12, color=DARK_TEXT, space_after=2)

# Right side - Pipeline & Industry-Academia
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(0.5),
             "Research-to-Product Pipeline", font_size=18, color=NAVY, bold=True)

pipeline_stages = [
    ("Research", "Perception model\n& evaluation framework"),
    ("Evaluation", "Subjective listening test\n& psychoacoustic analysis"),
    ("Validation", "System-level verification\n& regulatory compliance"),
    ("Production", "Mass production tuning\n& quality management"),
]

for i, (stage, desc) in enumerate(pipeline_stages):
    y = Inches(2.0) + Inches(i * 1.15)
    # Stage label
    label_box = add_shape(slide, Inches(7.2), y, Inches(1.8), Inches(0.8), fill_color=NAVY)
    tf_l = label_box.text_frame
    tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_l.paragraphs[0].text = stage
    tf_l.paragraphs[0].font.size = Pt(13)
    tf_l.paragraphs[0].font.color.rgb = WHITE
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Desc
    add_text_box(slide, Inches(9.2), y, Inches(3.3), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY, line_spacing=1.3)
    # Arrow
    if i < 3:
        add_text_box(slide, Inches(7.9), y + Inches(0.8), Inches(0.4), Inches(0.35),
                     "↓", font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Industry-academia joint
add_text_box(slide, Inches(7.2), Inches(6.0), Inches(5.3), Inches(0.5),
             "Industry-Academia Joint Research (as PI)", font_size=14, color=NAVY, bold=True)
add_text_box(slide, Inches(7.2), Inches(6.4), Inches(5.3), Inches(0.6),
             "• Seoul National Univ.: Consumer-driven road noise research\n"
             "• Chungnam National Univ.: AVAS soundscape evaluation (NRF funded)",
             font_size=12, color=MID_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════
# S10. IMPACT 4: HEALTH
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Impact 4: Health — Sound for Human Well-being")
add_slide_number(slide, 10)

# AI Diagnosis
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "AI-based Respiratory Disease Diagnosis", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(1.75))

health_ai = [
    "• Cough sound classification using LSTM/CNN",
    "  → Pneumonia diagnosis accuracy: 84.9%",
    "  → Chung, Jo et al. (2021) Sensors, SCI-E",
    "",
    "• Room Impulse Response (RIR) data augmentation",
    "  → Deep learning sound classifier performance improvement",
    "  → Patent registered (Korea + US Patent No. 18273592)",
    "  → Technology transfer: $50,000 to Hanyang S&A",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(3.0),
                  health_ai[0], font_size=13, color=DARK_TEXT, line_spacing=1.3)
for item in health_ai[1:]:
    c = ACCENT_BLUE if "→" in item else DARK_TEXT
    add_paragraph(tf, item, font_size=13, color=c, space_after=2)

# Digital Therapeutics
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.5),
             "VR Soundscape Digital Therapeutics", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(5.15))

health_vr = [
    "• VR-based soundscape therapy for mental illness",
    "  (Depression, schizophrenia — Hanyang Univ. Hospital)",
    "",
    "• Psycho-physiological restoration quantified",
    "  via EEG, HRV, and eye-tracking",
    "  → Jeon, Jo* & Lee (2023) Sustainable Cities & Society, Q1",
    "  → Jo et al. (2022) Scientific Report, Q1",
]
tf2 = add_text_box(slide, Inches(0.8), Inches(5.35), Inches(5.5), Inches(2.0),
                   health_vr[0], font_size=13, color=DARK_TEXT, line_spacing=1.3)
for item in health_vr[1:]:
    c = ACCENT_BLUE if "→" in item else DARK_TEXT
    add_paragraph(tf2, item, font_size=13, color=c, space_after=2)

# Right - Key message
msg_box = add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
                    fill_color=LIGHT_BLUE_BG)
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.0), Inches(0.5),
             "Expanding Building Acoustics", font_size=16, color=NAVY, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5.0), Inches(4.0),
             "Building acoustics methodologies\n— sound analysis, spatial audio,\npsychoacoustic evaluation —\n\n"
             "can be extended to:\n\n"
             "✓ Healthcare (disease diagnosis)\n"
             "✓ Mental health (digital therapeutics)\n"
             "✓ Wellness (restorative environments)\n\n"
             "This demonstrates the broad applicability\n"
             "of perception-driven sound research\n"
             "beyond traditional building engineering.",
             font_size=14, color=MID_GRAY, line_spacing=1.35)


# ═══════════════════════════════════════════════════
# S11. GLOBAL NETWORK
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Global Network & Scholarly Activities")
add_slide_number(slide, 11)

# Left - Collaborations
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "International Joint Research", font_size=18, color=NAVY, bold=True)

joint_items = [
    ("SATP Project (2020-2022)", "Soundscape Attributes Translation — 18 languages\nUCL, Stockholm Univ., TU Berlin, McGill, NUS, etc."),
    ("STAR Project (2018-2020)", "France-Korea urban soundscape comparison\nSorbonne University, Paris"),
    ("CHIC Project (2016-2019)", "VR 3D audio environment modeling\nRWTH Aachen University, KIST"),
]
for i, (title, desc) in enumerate(joint_items):
    y = Inches(1.9) + Inches(i * 1.5)
    add_text_box(slide, Inches(0.8), y, Inches(5.5), Inches(0.4),
                 title, font_size=14, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.35), Inches(5.5), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY, line_spacing=1.3)

# Right - Scholarly
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(0.5),
             "Scholarly Activities", font_size=18, color=NAVY, bold=True)

scholarly_items = [
    "Journal Reviewer: 17 international journals",
    "  Including B&E, L&UP, SCS, Applied Acoustics,",
    "  Science of Total Environment, IEEE Access, etc.",
    "",
    "Keynote Speaker:",
    "  Urban Sound Symposium, E-congress (2021)",
    "  \"Public space soundscapes — three continents\"",
    "",
    "Professional Affiliations:",
    "  ASA | AES | KSNVE | ASK | AIK | SAREK",
    "",
    "Awards (selected):",
    "  • EAA Best Paper & Presentation (ICA 2019)",
    "  • I-INCE Young Professional Award (2020)",
    "  • HMG Special Award for Research (2025)",
    "  • Promising Scientist Award, KSNVE (2022)",
]
tf = add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(5.0),
                  scholarly_items[0], font_size=13, color=DARK_TEXT, bold=True, line_spacing=1.25)
for item in scholarly_items[1:]:
    b = bool(not item.startswith("  ") and item != "" and not item.startswith("•"))
    add_paragraph(tf, item, font_size=12 if item.startswith("  ") else 13,
                  color=DARK_TEXT, bold=b, space_after=1)


# ═══════════════════════════════════════════════════
# S12. WHAT INDUSTRY TAUGHT ME
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "What Industry Taught Me")
add_slide_number(slide, 12)

lessons = [
    ("Research-to-Product Pipeline",
     "Experienced the full cycle: perception model → listening test\n"
     "→ system verification → mass production tuning.\n\n"
     "Research outcomes were validated not only through\npeer review, but through real-world product deployment."),
    ("Industry-Ready Research Network",
     "Managed industry-academia joint research as PI\n"
     "from the corporate side (Seoul Nat'l Univ., Chungnam Nat'l Univ.).\n\n"
     "Ready to initiate joint research with Hyundai Motor, Samsung,\n"
     "and construction companies from Day 1."),
    ("Real-World Problem Definition",
     "Industry problems are the best research questions.\n\n"
     "\"How should AVAS sound in 2030 cities?\" — this question\n"
     "came from the factory floor and became a JASA paper."),
]

for i, (title, desc) in enumerate(lessons):
    x = Inches(0.6) + Inches(i * 4.2)
    y = Inches(1.5)

    box = add_shape(slide, x, y, Inches(3.9), Inches(4.8), fill_color=LIGHT_BLUE_BG)

    # Number
    num_box = add_shape(slide, x + Inches(0.2), y + Inches(0.2),
                         Inches(0.5), Inches(0.5), fill_color=NAVY)
    tf_n = num_box.text_frame
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.paragraphs[0].text = str(i + 1)
    tf_n.paragraphs[0].font.size = Pt(18)
    tf_n.paragraphs[0].font.color.rgb = WHITE
    tf_n.paragraphs[0].font.bold = True
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.8), Inches(0.6),
                 title, font_size=16, color=NAVY, bold=True)

    # Desc
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.3), Inches(3.5),
                 desc, font_size=13, color=MID_GRAY, line_spacing=1.4)

# Bottom quote
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
             "\"This experience shaped how I define, validate, and deliver research.\"",
             font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# S13. VISION — SENSE Lab
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_slide_number(slide, 13, 20)

# Section header
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.5),
             "FUTURE RESEARCH", font_size=16, color=ACCENT_TEAL, bold=True)

# Lab name
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
             "SENSE Lab", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.6),
             "Sound Environment aNd Sensory Engineering Laboratory",
             font_size=22, color=ACCENT_TEAL)

# Vision statement
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.8),
             "Completing the building environment research at Hanyang:",
             font_size=18, color=WHITE)

# Three pillars diagram
pillars = [
    ("Thermal\nEnvironment", "(Existing)", RGBColor(0x22, 0x44, 0x66)),
    ("Air Quality\nEnvironment", "(Existing)", RGBColor(0x22, 0x44, 0x66)),
    ("Sound\nEnvironment", "(SENSE Lab)", ACCENT_BLUE),
]
for i, (name, status, color) in enumerate(pillars):
    x = Inches(0.8) + Inches(i * 3.5)
    y = Inches(4.5)
    box = add_shape(slide, x, y, Inches(3.0), Inches(1.3), fill_color=color)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = status
    p2.font.size = Pt(12)
    p2.font.color.rgb = ACCENT_TEAL
    p2.alignment = PP_ALIGN.CENTER

# Arrow down to integration
add_text_box(slide, Inches(5.0), Inches(5.85), Inches(3.3), Inches(0.4),
             "▼              ▼              ▼", font_size=14, color=ACCENT_TEAL,
             alignment=PP_ALIGN.CENTER)

# Integration box
int_box = add_shape(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.9),
                    fill_color=ACCENT_BLUE)
tf_int = int_box.text_frame
tf_int.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_int.paragraphs[0].text = "Intelligent Building Environment System (IBES)"
tf_int.paragraphs[0].font.size = Pt(20)
tf_int.paragraphs[0].font.color.rgb = WHITE
tf_int.paragraphs[0].font.bold = True
tf_int.paragraphs[0].alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════
# S14. PHASE 1 + PHASE 2
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Research Plan — Phase 1 & Phase 2")
add_slide_number(slide, 14)

# Phase 1
p1_box = add_shape(slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.8),
                   fill_color=LIGHT_BLUE_BG)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(2.5), Inches(0.4),
             "PHASE 1  |  Year 1-3", font_size=15, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.85), Inches(5.5), Inches(0.6),
             "AI-driven Building Acoustic\nEnvironment Diagnosis",
             font_size=20, color=NAVY, bold=True, line_spacing=1.1)

phase1_items = [
    "① IoT acoustic sensor network for indoor/outdoor",
    "    sound environment monitoring",
    "",
    "② ML-based sound environment auto-classification",
    "    and anomaly detection",
    "",
    "③ VR sound environment simulation platform",
    "    for design-stage acoustic previewing",
    "",
    "Target Funding:",
    "  • NRF Young/Mid-career Researcher (₩200-300M/yr)",
    "  • MSIT Smart City Program",
    "  • Hyundai Motor joint research (AVAS/NVH)",
    "",
    "Expected: 3-4 SCI papers/yr, 3-4 graduate students",
]
tf = add_text_box(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(4.0),
                  phase1_items[0], font_size=12, color=DARK_TEXT, line_spacing=1.25)
for item in phase1_items[1:]:
    b = bool(item.startswith("Target") or item.startswith("Expected"))
    c = NAVY if b else (ACCENT_BLUE if item.startswith("  •") else DARK_TEXT)
    add_paragraph(tf, item, font_size=12, color=c, bold=b, space_after=1)

# Phase 2
p2_box = add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.8),
                   fill_color=LIGHT_BLUE_BG)
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(2.5), Inches(0.4),
             "PHASE 2  |  Year 3-5", font_size=15, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.6),
             "Mobility-Urban Sound\nEnvironment Interaction",
             font_size=20, color=NAVY, bold=True, line_spacing=1.1)

phase2_items = [
    "① Long-term monitoring of urban sound environment",
    "    changes in the EV transition era",
    "",
    "② AVAS + urban soundscape integrated design",
    "    framework and guidelines",
    "",
    "③ Digital twin-based sound environment",
    "    prediction and simulation",
    "",
    "Target Funding:",
    "  • NRF Leading Research Center (₩500-1000M/yr)",
    "  • MOLIT / MOE urban noise management",
    "  • Samsung — Spatial Audio / immersive media",
    "",
    "Expected: ISO standard participation, 6-8 students",
]
tf2 = add_text_box(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(4.0),
                   phase2_items[0], font_size=12, color=DARK_TEXT, line_spacing=1.25)
for item in phase2_items[1:]:
    b = bool(item.startswith("Target") or item.startswith("Expected"))
    c = NAVY if b else (ACCENT_BLUE if item.startswith("  •") else DARK_TEXT)
    add_paragraph(tf2, item, font_size=12, color=c, bold=b, space_after=1)


# ═══════════════════════════════════════════════════
# S15. PHASE 3 + INTEGRATION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Research Plan — Phase 3 & Integration Vision")
add_slide_number(slide, 15)

# Phase 3
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.4),
             "PHASE 3  |  Year 5+", font_size=15, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(6), Inches(0.6),
             "Multi-modal Intelligent Building Environment Platform",
             font_size=20, color=NAVY, bold=True)

phase3_items = [
    "① Integrated sensing & control: Sound + Thermal + Air Quality",
    "    → Joint research with existing faculty at Hanyang AE dept.",
    "",
    "② Human-centered building environment optimization",
    "    → Wellness, productivity, and energy efficiency simultaneously",
    "",
    "③ International research hub leveraging existing network",
    "    → UCL, RWTH Aachen, Sorbonne partnerships",
    "    → ISO 12913 standard expansion and leadership",
]
tf = add_text_box(slide, Inches(0.8), Inches(2.3), Inches(6.5), Inches(3.5),
                  phase3_items[0], font_size=13, color=DARK_TEXT, line_spacing=1.3)
for item in phase3_items[1:]:
    c = ACCENT_BLUE if "→" in item else DARK_TEXT
    add_paragraph(tf, item, font_size=13, color=c, space_after=2)

# Funding roadmap (right side)
add_text_box(slide, Inches(7.5), Inches(1.3), Inches(5), Inches(0.5),
             "Funding Roadmap", font_size=18, color=NAVY, bold=True)

roadmap = [
    ("Year 1-2", "~₩400M/yr", "NRF Young Researcher\n+ Hyundai Motor joint research"),
    ("Year 3-4", "~₩700M/yr", "NRF Mid-career + MOLIT\n+ Industry partners (2-3)"),
    ("Year 5+", "~₩1B+/yr", "Leading Research Center\n+ International grants"),
]
for i, (year, amount, desc) in enumerate(roadmap):
    y = Inches(2.0) + Inches(i * 1.5)
    box = add_shape(slide, Inches(7.5), y, Inches(5.0), Inches(1.2), fill_color=LIGHT_BLUE_BG)
    add_text_box(slide, Inches(7.7), y + Inches(0.1), Inches(1.5), Inches(0.4),
                 year, font_size=14, color=NAVY, bold=True)
    add_text_box(slide, Inches(9.3), y + Inches(0.1), Inches(1.5), Inches(0.4),
                 amount, font_size=14, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(7.7), y + Inches(0.5), Inches(4.5), Inches(0.7),
                 desc, font_size=12, color=MID_GRAY, line_spacing=1.3)

# Ready-to-launch
add_text_box(slide, Inches(7.5), Inches(6.0), Inches(5), Inches(0.4),
             "Ready to Launch from Day 1:", font_size=14, color=NAVY, bold=True)
add_text_box(slide, Inches(7.5), Inches(6.35), Inches(5), Inches(0.8),
             "• Hyundai Motor — existing relationship, AVAS/NVH\n"
             "• Samsung — spatial audio & immersive media\n"
             "• NRF — 11 funded projects on record\n"
             "• UCL / RWTH / Sorbonne — active collaborations",
             font_size=12, color=MID_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════
# S16. FUNDING DETAIL (expanded from previous roadmap)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Funding Strategy — Detailed Plan")
add_slide_number(slide, 16)

# Government
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "Government & Public Funding", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(1.75))

gov_items = [
    "National Research Foundation (NRF):",
    "  • Young/Mid-career Researcher: Soundscape AI (₩100-200M/yr)",
    "  • International Joint: UCL (UK), Sorbonne (France) (₩120M/yr)",
    "  • Basic Research Lab: Building Env. Integration (₩500M/3yr)",
    "",
    "Government Agencies:",
    "  • MSIT: Smart city sound environment monitoring",
    "  • MOLIT/KAIA: Urban noise management & policy",
    "  • MOE: Environmental noise impact assessment",
    "  • LH Corporation: Residential noise quality standards",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(3.5),
                  gov_items[0], font_size=13, color=NAVY, bold=True, line_spacing=1.25)
for item in gov_items[1:]:
    b = bool(not item.startswith("  ") and item != "")
    add_paragraph(tf, item, font_size=12 if item.startswith("  ") else 13,
                  color=NAVY if b else DARK_TEXT, bold=b, space_after=2)

# Industry
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(0.5),
             "Industry Collaboration", font_size=18, color=NAVY, bold=True)

industry_items = [
    ("Hyundai Motor Company", "AVAS sound design, vehicle cabin acoustics,\nconsumer satisfaction modeling\n(Existing relationship — joint research since 2023)"),
    ("Samsung Research", "Spatial audio quality evaluation,\nimmersive media perception modeling\n(Demonstrated capability in portfolio)"),
    ("Construction / Architecture", "Building acoustics consulting,\nnoise impact assessment, office acoustics\n(FURSYS precedent: ₩60M project)"),
]
for i, (company, desc) in enumerate(industry_items):
    y = Inches(2.0) + Inches(i * 1.7)
    box = add_shape(slide, Inches(7.2), y, Inches(5.3), Inches(1.4), fill_color=LIGHT_BLUE_BG)
    add_text_box(slide, Inches(7.5), y + Inches(0.1), Inches(4.8), Inches(0.4),
                 company, font_size=14, color=NAVY, bold=True)
    add_text_box(slide, Inches(7.5), y + Inches(0.45), Inches(4.8), Inches(0.9),
                 desc, font_size=11, color=MID_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════
# S17. TEACHING PHILOSOPHY
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Teaching Philosophy")
add_slide_number(slide, 17)

# Vision
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.8),
             "\"Educating future-ready architectural engineers who understand\nboth technology and human experience\"",
             font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER,
             line_spacing=1.3)

# Three principles
teach_pillars = [
    ("Theory → Practice",
     "Connect classroom theory to real-world\napplications through industry case studies.\n\n"
     "Use actual projects from Hyundai Motor,\nconstruction sites, and concert hall designs\n"
     "as teaching materials.\n\n"
     "VR/AR simulation labs for immersive\nlearning experiences."),
    ("Data Literacy",
     "Equip architectural engineers with\nAI and data analysis capabilities.\n\n"
     "Python/MATLAB-based acoustic data\nprocessing and visualization.\n\n"
     "Machine learning fundamentals applied\n"
     "to building performance optimization\nand environmental monitoring."),
    ("Global Exposure",
     "Mandatory international conference\npresentation for graduate students.\n\n"
     "Short-term research visits to partner\ninstitutions (UCL, RWTH, Sorbonne).\n\n"
     "Bilingual research training to prepare\nstudents for global careers in both\nacademia and industry."),
]

for i, (title, desc) in enumerate(teach_pillars):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.5)
    box = add_shape(slide, x, y, Inches(3.8), Inches(4.5), fill_color=LIGHT_BLUE_BG)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                 title, font_size=18, color=NAVY, bold=True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Inches(0.3), y + Inches(0.8), Inches(1.5), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(3.2),
                 desc, font_size=13, color=MID_GRAY, line_spacing=1.35)


# ═══════════════════════════════════════════════════
# S18. TEACHING PLAN + MENTORING
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Teaching Plan & Student Mentoring")
add_slide_number(slide, 18)

# Courses - Left
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "Course Portfolio", font_size=18, color=NAVY, bold=True)
add_accent_bar(slide, Inches(1.75))

# Existing courses
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
             "Courses to Inherit:", font_size=14, color=ACCENT_BLUE, bold=True)
courses_exist = [
    "• Building Environmental Engineering (건축환경공학) — UG",
    "• Architectural Acoustics (건축음향) — UG",
    "• Probability & Statistics (확률통계론) — UG",
]
tf = add_text_box(slide, Inches(0.8), Inches(2.35), Inches(5.5), Inches(1.2),
                  courses_exist[0], font_size=13, color=DARK_TEXT, line_spacing=1.3)
for c in courses_exist[1:]:
    add_paragraph(tf, c, font_size=13, color=DARK_TEXT, space_after=3)

add_text_box(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.4),
             "New Courses to Develop:", font_size=14, color=ACCENT_BLUE, bold=True)
courses_new = [
    "• Smart Building Environment Systems — Grad",
    "  (IoT + ML + VR for building environment design)",
    "• AI for Built Environment — Grad",
    "  (Python/MATLAB acoustic data analysis lab)",
    "• Sustainable Building & Ecology — UG",
    "  (LEED/WELL standards, green building practice)",
]
tf2 = add_text_box(slide, Inches(0.8), Inches(3.85), Inches(5.5), Inches(2.5),
                   courses_new[0], font_size=13, color=DARK_TEXT, line_spacing=1.25)
for c in courses_new[1:]:
    sz = 12 if c.startswith("  ") else 13
    add_paragraph(tf2, c, font_size=sz, color=MID_GRAY if c.startswith("  ") else DARK_TEXT, space_after=2)

# Right - Mentoring
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(0.5),
             "Student Mentoring", font_size=18, color=NAVY, bold=True)

# Undergrad
add_text_box(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(0.4),
             "Undergraduate:", font_size=14, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.2), Inches(2.35), Inches(5.3), Inches(1.0),
             "• Capstone design linked to industry projects\n"
             "• VR/BIM hands-on lab experiences\n"
             "• Early research exposure program",
             font_size=13, color=DARK_TEXT, line_spacing=1.3)

# Grad
add_text_box(slide, Inches(7.2), Inches(3.5), Inches(5.3), Inches(0.4),
             "Graduate:", font_size=14, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.2), Inches(3.85), Inches(5.3), Inches(1.2),
             "• Target: 1 SCI paper + 1 patent per student\n"
             "• International conference presentation required\n"
             "• Short-term visit to UCL/RWTH/Sorbonne\n"
             "• Weekly lab seminar + writing workshop",
             font_size=13, color=DARK_TEXT, line_spacing=1.3)

# Career
add_text_box(slide, Inches(7.2), Inches(5.2), Inches(5.3), Inches(0.4),
             "Career Development:", font_size=14, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.2), Inches(5.55), Inches(5.3), Inches(1.2),
             "• Hyundai Motor / Kia — NVH team\n"
             "• Samsung Research — Audio Lab\n"
             "• Construction & architecture firms\n"
             "• Government research institutes (KICT, KRICT)",
             font_size=13, color=DARK_TEXT, line_spacing=1.3)


# ═══════════════════════════════════════════════════
# S19. CONTRIBUTION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_content_slide(slide, "Contribution to the Department")
add_slide_number(slide, 19)

contributions = [
    ("Completing Building\nEnvironment Research",
     "Sound joins Thermal and Air Quality to form\na comprehensive building environment research group.\n\n"
     "Enables joint BK21 / Leading Research Center\napplications with existing faculty.",
     "01"),
    ("Expanding Industry\nCollaboration",
     "Hyundai Motor and Samsung networks\nopen new industry partnership channels\nfor the entire department.\n\n"
     "Construction and architecture firm consulting\nthrough acoustics expertise.",
     "02"),
    ("Elevating Global\nPresence",
     "Active collaborations with UCL, RWTH Aachen,\nand Sorbonne University.\n\n"
     "Potential for department-level MOU\nand student exchange programs.\n\n"
     "ISO 12913 standardization leadership.",
     "03"),
    ("Hanyang Alumni\nCommitment",
     "B.S. and Ph.D. from Hanyang — deep understanding\nof department culture and history.\n\n"
     "Continuing the tradition of building acoustics\nresearch established at this department.\n\n"
     "Ready to contribute to department governance\nand student mentoring from Day 1.",
     "04"),
]

for i, (title, desc, num) in enumerate(contributions):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(1.3) + Inches(row * 3.0)

    box = add_shape(slide, x, y, Inches(6.0), Inches(2.7), fill_color=LIGHT_BLUE_BG)

    # Number
    num_box = add_shape(slide, x + Inches(0.2), y + Inches(0.2),
                         Inches(0.5), Inches(0.5), fill_color=NAVY)
    tf_n = num_box.text_frame
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.paragraphs[0].text = num
    tf_n.paragraphs[0].font.size = Pt(16)
    tf_n.paragraphs[0].font.color.rgb = WHITE
    tf_n.paragraphs[0].font.bold = True
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, x + Inches(0.9), y + Inches(0.15), Inches(4.8), Inches(0.7),
                 title, font_size=16, color=NAVY, bold=True, line_spacing=1.1)

    # Desc
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(5.4), Inches(1.7),
                 desc, font_size=12, color=MID_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════
# S20. THANK YOU
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

# Quote
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.5),
             "From Hanyang's tradition\nin building acoustics,\nthrough industry validation,\n"
             "to intelligent sound environments\n— SENSE Lab is ready.",
             font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             line_spacing=1.4)

# Separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(5.5), Inches(4.3), Inches(2.3), Pt(2))
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# Name & Contact
add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.6),
             "Hyun In Jo, Ph.D.", font_size=28, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.4),
             "best2012@naver.com  |  linkedin.com/in/hyunin-jo",
             font_size=14, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
             "Department of Architectural Engineering\nHanyang University",
             font_size=16, color=RGBColor(0x88, 0x99, 0xAA),
             alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Thank you
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
             "THANK YOU", font_size=18, color=ACCENT_TEAL, bold=True,
             alignment=PP_ALIGN.CENTER)


# ─── Save ───
output_path = "/Users/hyunbin/Research/HanyangUniv_Faculty_Presentation_HyunInJo.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
