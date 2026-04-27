#!/usr/bin/env python3
"""
Hanyang University Faculty Interview Presentation v2
- Samsung Research color scheme
- Full paper citations as footnotes
- Bilingual speaker notes (English + Korean)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Samsung Research Exact Color Palette ───
SAMSUNG_DARK = RGBColor(0x0A, 0x0F, 0x1E)       # Near-black navy (slide BG)
SAMSUNG_BLUE = RGBColor(0x13, 0x28, 0x9F)       # #13289F Samsung Research primary
SAMSUNG_BLUE2 = RGBColor(0x3D, 0x7D, 0xDE)      # #3D7DDE secondary blue
SAMSUNG_LIGHT_BLUE = RGBColor(0x06, 0x89, 0xD8) # #0689D8 accent/arrow blue
SAMSUNG_CYAN = RGBColor(0x3D, 0x7D, 0xDE)       # Same as BLUE2 for accents
GRAY_BLUE = RGBColor(0x5B, 0x71, 0x8D)          # #5B718D tertiary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xE7, 0xE6, 0xE2)          # #E7E6E2
LIGHT_CARD = RGBColor(0xF0, 0xF4, 0xF8)         # #F0F4F8
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)           # #1A1A2E
MID_GRAY = RGBColor(0x75, 0x78, 0x7B)           # #75787B
FOOTNOTE_GRAY = RGBColor(0x88, 0x88, 0x99)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
SEPARATOR = RGBColor(0x22, 0x33, 0x55)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ─── Helpers ───
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.line.fill.background()
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    return s

def add_rounded_shape(slide, left, top, width, height, fill_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.line.fill.background()
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    return s

def tb(slide, left, top, width, height, text, sz=18, color=DARK_TEXT,
       bold=False, align=PP_ALIGN.LEFT, font="Calibri", lsp=1.2):
    """Add a text box — returns the text_frame."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(sz * (lsp - 1))
    return tf

def ap(tf, text, sz=16, color=DARK_TEXT, bold=False,
       align=PP_ALIGN.LEFT, sa=4, sb=0, level=0):
    """Add a paragraph to existing text_frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bool(bold)
    p.font.name = "Calibri"
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    p.level = level
    return p

def accent_bar(slide, top, left=Inches(0.8), width=Inches(1.5), color=SAMSUNG_BLUE):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()

def content_header(slide, title):
    """White slide with Samsung-style header."""
    add_bg(slide, WHITE)
    # Top bar
    bar = add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), SAMSUNG_BLUE)
    # Title
    tb(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.75),
       title, sz=26, color=SAMSUNG_DARK, bold=True)
    # Line under title
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(11.7), Pt(1.5), SAMSUNG_BLUE)

def slide_num(slide, n, total=20):
    tb(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
       f"{n} / {total}", sz=9, color=FOOTNOTE_GRAY, align=PP_ALIGN.RIGHT)

def footnote(slide, text, y=Inches(6.35)):
    """Add citation footnote at bottom of slide."""
    tb(slide, Inches(0.5), y, Inches(12.5), Inches(1.1),
       text, sz=7, color=FOOTNOTE_GRAY, lsp=1.15)

def add_notes(slide, en_text, kr_text):
    """Add bilingual speaker notes."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "[English Script]"
    p.font.bold = True
    p.font.size = Pt(11)
    for line in en_text.strip().split("\n"):
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
    # Korean
    p3 = tf.add_paragraph()
    p3.text = ""
    p4 = tf.add_paragraph()
    p4.text = "[한국어 스크립트]"
    p4.font.bold = True
    p4.font.size = Pt(11)
    for line in kr_text.strip().split("\n"):
        p5 = tf.add_paragraph()
        p5.text = line
        p5.font.size = Pt(11)


# ═══════════════════════════════════════
# S1. TITLE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SAMSUNG_DARK)

accent_bar(s, Inches(1.8), Inches(0.8), Inches(2), SAMSUNG_LIGHT_BLUE)

tb(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
   "Intelligent Sound Environment Systems\nfor Human-Centered Built Environments",
   sz=36, color=WHITE, bold=True, lsp=1.3)

tb(s, Inches(0.8), Inches(3.8), Inches(8), Inches(0.5),
   "Faculty Position: Intelligent Building Environment Systems",
   sz=17, color=SAMSUNG_LIGHT_BLUE)

add_shape(s, Inches(0.8), Inches(4.6), Inches(5), Pt(1), SEPARATOR)

tb(s, Inches(0.8), Inches(4.9), Inches(6), Inches(0.6),
   "Hyun In Jo, Ph.D.", sz=28, color=WHITE, bold=True)

tb(s, Inches(0.8), Inches(5.5), Inches(8), Inches(0.35),
   "LEED AP BD+C  |  LEED AP ID+C  |  WELL AP  |  ADsP",
   sz=13, color=SAMSUNG_CYAN)

tb(s, Inches(0.8), Inches(6.0), Inches(10), Inches(0.5),
   "Department of Architectural Engineering, Hanyang University",
   sz=15, color=RGBColor(0x88, 0x99, 0xAA))

add_notes(s,
"""Good morning/afternoon. Thank you for the opportunity to present my research and vision.
My name is Hyun In Jo. I received both my Bachelor's and Ph.D. degrees from Hanyang University's Department of Architectural Engineering. After working as a post-doctoral researcher at KICT and a senior research engineer at Hyundai Motor Company for four years, I am now applying for the faculty position in Intelligent Building Environment Systems.
Today, I will present my research achievements, future research plan, and how I can contribute to this department.""",
"""안녕하십니까. 발표 기회를 주셔서 감사합니다.
저는 한양대학교 건축공학과에서 학사와 박사 학위를 취득한 조현인입니다. 한국건설기술연구원 박사후연구원과 현대자동차 연구소에서 4년간 책임연구원으로 근무한 뒤, 지능형 건축환경 시스템 분야 전임교원에 지원하였습니다.
오늘 저의 연구 실적, 향후 연구 계획, 그리고 학과에 대한 기여 방안을 말씀드리겠습니다.""")


# ═══════════════════════════════════════
# S2. TABLE OF CONTENTS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Table of Contents")
slide_num(s, 2)

sections = [
    ("01", "Introduction", "The Evolution  |  Brief Introduction"),
    ("02", "Research", "Philosophy  |  Key Research (Buildings · Cities · Products · Health)\nGlobal Network  |  Industry Experience"),
    ("03", "Future Research Plan", "SENSE Lab Vision  |  Phase 1-2-3  |  Funding Roadmap"),
    ("04", "Education", "Teaching Philosophy  |  Teaching Plan  |  Student Mentoring"),
    ("05", "Contribution", "Department  |  Industry  |  Global Network"),
]
for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.4) + Inches(i * 1.1)
    box = add_rounded_shape(s, Inches(0.8), y, Inches(0.8), Inches(0.65), SAMSUNG_DARK)
    tf_n = box.text_frame
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.paragraphs[0].text = num
    tf_n.paragraphs[0].font.size = Pt(18)
    tf_n.paragraphs[0].font.color.rgb = WHITE
    tf_n.paragraphs[0].font.bold = True
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb(s, Inches(1.8), y + Inches(0.05), Inches(3), Inches(0.6),
       title, sz=20, color=SAMSUNG_DARK, bold=True)
    tb(s, Inches(4.8), y + Inches(0.05), Inches(7.5), Inches(0.6),
       desc, sz=13, color=MID_GRAY, lsp=1.3)

add_notes(s,
"""Here is the outline of today's presentation. I will start with a brief introduction, then discuss my research achievements across four impact areas, followed by my future research plan under SENSE Lab, and finally my education and contribution plans.""",
"""오늘 발표의 목차입니다. 간략한 자기소개 후, 네 가지 영향 영역에 걸친 연구 실적을 말씀드리고, SENSE Lab으로서의 향후 연구 계획, 마지막으로 교육 및 기여 계획을 발표하겠습니다.""")


# ═══════════════════════════════════════
# S3. THE EVOLUTION
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "The Evolution of Building Sound Environment Research")
slide_num(s, 3)

gen_data = [
    ("Generation 1", "Noise Control",
     "Physical metrics\ndB, RT60, NC curves\n\n\"How loud is it?\"",
     RGBColor(0xE5, 0xE8, 0xED), DARK_TEXT),
    ("Generation 2", "Soundscape",
     "Perceptual experience\nISO 12913, VR, Bio-signals\n\n\"How do people\nexperience it?\"",
     RGBColor(0xD0, 0xE4, 0xF5), SAMSUNG_DARK),
    ("Generation 3", "Intelligent\nSound Environment",
     "Data-driven optimization\nAI, IoT, Digital Twin\n\n\"How to optimize\nautomatically?\"",
     SAMSUNG_DARK, WHITE),
]
for i, (gen, title, desc, bg, tc) in enumerate(gen_data):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.4)
    add_rounded_shape(s, x, y, Inches(3.7), Inches(4.6), bg)
    tb(s, x + Inches(0.3), y + Inches(0.2), Inches(3.1), Inches(0.35),
       gen, sz=12, color=SAMSUNG_BLUE if i < 2 else SAMSUNG_CYAN, bold=True)
    tb(s, x + Inches(0.3), y + Inches(0.6), Inches(3.1), Inches(1.2),
       title, sz=24, color=tc, bold=True, lsp=1.1)
    tb(s, x + Inches(0.3), y + Inches(1.9), Inches(3.1), Inches(2.2),
       desc, sz=14, color=tc if i == 2 else MID_GRAY, lsp=1.35)
    if i < 2:
        tb(s, x + Inches(3.75), Inches(3.2), Inches(0.3), Inches(0.5),
           "→", sz=24, color=SAMSUNG_BLUE, bold=True, align=PP_ALIGN.CENTER)

tb(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
   "\"I have researched across all three generations — and I am ready to lead Generation 3.\"",
   sz=17, color=SAMSUNG_DARK, bold=True, align=PP_ALIGN.CENTER)

add_notes(s,
"""The field of building sound environment research has evolved through three generations.
Generation 1 was traditional noise control — focused on physical metrics like decibels and reverberation time. The question was simply "how loud is it?"
Generation 2 introduced the soundscape paradigm — shifting focus to human perception and experience, guided by ISO 12913 standards. The question became "how do people actually experience the sound environment?"
Generation 3, which I propose, is the intelligent sound environment — using AI, IoT sensors, and digital twins to automatically optimize building acoustic environments. The question becomes "how can we optimize environments automatically based on human needs?"
I have researched across all three generations and I am ready to lead Generation 3 at Hanyang University.""",
"""건축 음환경 연구는 세 세대를 거쳐 진화해 왔습니다.
1세대는 전통적 소음 제어입니다. 데시벨, 잔향시간 등 물리량 중심으로 "얼마나 시끄러운가?"가 핵심 질문이었습니다.
2세대는 사운드스케이프 패러다임으로, ISO 12913 표준에 기반하여 인간의 지각과 경험으로 초점이 전환되었습니다. "사람들은 음환경을 어떻게 경험하는가?"가 핵심 질문이 되었습니다.
제가 제안하는 3세대는 지능형 음환경으로, AI, IoT 센서, 디지털 트윈을 활용하여 건축 음환경을 자동으로 최적화합니다.
저는 세 세대의 연구를 모두 수행해 왔으며, 한양대학교에서 3세대를 이끌 준비가 되어 있습니다.""")


# ═══════════════════════════════════════
# S4. BRIEF INTRODUCTION
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Brief Introduction")
slide_num(s, 4)

tl = [
    ("2013 – 2016", "B.S. Architectural Engineering",
     "Hanyang University (Summa Cum Laude, 6-semester Early Graduation)", SAMSUNG_BLUE),
    ("2016 – 2022", "Ph.D. Architectural Environmental Engineering",
     "Hanyang University (GPA 4.39 / 4.5, Advisor: Prof. Jin Yong Jeon)", SAMSUNG_BLUE),
    ("2018.10", "Visiting Researcher",
     "Institut Jean Le Rond d'Alembert, Sorbonne University, Paris", SAMSUNG_CYAN),
    ("2022.03 – 08", "Post-doctoral Researcher (NST Young Scientist Fellow)",
     "Korea Institute of Civil Engineering and Building Technology (KICT)", SAMSUNG_CYAN),
    ("2022.09 – Present", "Senior Research Engineer",
     "Hyundai Motor Company, NVH Test Research Lab", GOLD),
]
for i, (period, role, org, dc) in enumerate(tl):
    y = Inches(1.3) + Inches(i * 0.82)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.05), Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = dc; dot.line.fill.background()
    if i < len(tl) - 1:
        vl = add_shape(s, Inches(1.055), y + Inches(0.2), Pt(2), Inches(0.65), RGBColor(0xCC, 0xCC, 0xCC))
    tb(s, Inches(1.35), y - Inches(0.05), Inches(2.0), Inches(0.3),
       period, sz=12, color=MID_GRAY, bold=True)
    tb(s, Inches(3.35), y - Inches(0.05), Inches(4.0), Inches(0.3),
       role, sz=13, color=SAMSUNG_DARK, bold=True)
    tb(s, Inches(3.35), y + Inches(0.2), Inches(4.5), Inches(0.3),
       org, sz=11, color=MID_GRAY)

# Stats
stats = [("26", "SCI(E) Papers\n(21 First Author)"), ("18", "h-index"),
         ("6", "Patents\n(incl. US)"), ("12", "Awards"),
         ("$5M+", "Funded\nResearch"), ("4,565", "Concert Hall\nSeats")]
for i, (num, label) in enumerate(stats):
    r, c = i // 3, i % 3
    x = Inches(8.5) + Inches(c * 1.5)
    y = Inches(1.4) + Inches(r * 2.3)
    tb(s, x, y, Inches(1.4), Inches(0.7), num, sz=28, color=SAMSUNG_BLUE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x, y + Inches(0.7), Inches(1.4), Inches(0.8), label, sz=9, color=MID_GRAY, align=PP_ALIGN.CENTER, lsp=1.2)

add_notes(s,
"""Let me briefly introduce myself. I completed both my Bachelor's and Ph.D. at Hanyang University's Department of Architectural Engineering, graduating with highest honors in just 6 semesters for my bachelor's degree.
During my Ph.D., I conducted research on VR-based soundscape evaluation and psycho-physiological responses, and had the opportunity to visit Sorbonne University in Paris for joint research.
After graduating, I worked as a post-doctoral researcher at KICT on building energy data analysis, then joined Hyundai Motor Company where I have been working for four years as a senior research engineer in the NVH division.
As you can see on the right, I have published 26 SCI papers with 21 as first author, hold 6 patents including a US patent, received 12 research awards, and participated in over 5 million dollars of funded research.""",
"""간략히 자기소개를 드리겠습니다. 저는 한양대학교 건축공학과에서 학사와 박사를 모두 마쳤으며, 학사는 6학기 수석 조기졸업하였습니다.
박사과정 중 VR 기반 사운드스케이프 평가와 심리생리학적 반응 연구를 수행하였고, 프랑스 소르본대학교에 방문연구원으로 공동연구를 진행하였습니다.
졸업 후 한국건설기술연구원에서 건물에너지 데이터 분석 박사후연구원을 거쳐, 현대자동차 연구소 NVH 시험팀에서 4년간 책임연구원으로 근무하고 있습니다.
오른쪽에 보시는 바와 같이, SCI 논문 26편(주저자 21편), 특허 6건(미국 포함), 수상 12건, 총 연구비 50억원 이상의 과제에 참여한 실적이 있습니다.""")


# ═══════════════════════════════════════
# S5. RESEARCH PHILOSOPHY
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Research Philosophy")
slide_num(s, 5)

tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.8),
   "\"Designing human-centric acoustic environments\nby integrating perception science, architecture, and intelligent technology\"",
   sz=19, color=SAMSUNG_DARK, bold=True, align=PP_ALIGN.CENTER, lsp=1.3)

pillars = [
    ("Perception over Physics",
     "Evaluate environments through\nhuman experience, not just dB.\n\n"
     "Psychoacoustics, bio-signals\n(EEG, HRV, eye-tracking), and\n"
     "subjective assessment replace\npurely physical metrics."),
    ("Interdisciplinary Integration",
     "Bridge architectural acoustics with\ncognitive psychology, AI / data science,\n"
     "and urban design.\n\n"
     "One core methodology applied\nacross buildings, cities, products,\nand healthcare."),
    ("Lab to Life",
     "Research must reach real buildings,\ncities, and products.\n\n"
     "3 concert halls designed,\nEV sound in mass production,\n"
     "digital therapeutics validated\nin clinical settings."),
]
for i, (title, desc) in enumerate(pillars):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.4)
    add_rounded_shape(s, x, y, Inches(3.8), Inches(4.4), LIGHT_CARD)
    tb(s, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
       title, sz=17, color=SAMSUNG_DARK, bold=True)
    accent_bar(s, y + Inches(0.75), x + Inches(0.3), Inches(1.5), SAMSUNG_BLUE)
    tb(s, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(3.0),
       desc, sz=13, color=MID_GRAY, lsp=1.4)

add_notes(s,
"""My research philosophy rests on three core principles.
First, Perception over Physics. I believe we should evaluate building environments through how people actually experience them, not just through physical measurements. This is why I use psychoacoustic methods and bio-signals like EEG and eye-tracking alongside traditional acoustic metrics.
Second, Interdisciplinary Integration. My research bridges architectural acoustics with cognitive psychology, AI, and urban design. This single integrated methodology applies across buildings, cities, products, and even healthcare.
Third, Lab to Life. Research should not stay in journals. My work has been applied to three concert hall designs totaling over 4,500 seats, electric vehicle sounds in mass production, and VR-based digital therapeutics validated in clinical settings.""",
"""저의 연구 철학은 세 가지 핵심 원칙에 기반합니다.
첫째, 물리량보다 지각입니다. 건축환경을 물리적 측정값이 아닌, 사람이 실제로 경험하는 방식으로 평가해야 합니다. 이를 위해 심리음향 방법론과 뇌파, 시선추적 등 생체신호를 활용합니다.
둘째, 학제간 융합입니다. 건축음향을 인지심리학, AI, 도시설계와 연결하는 하나의 통합 방법론을 구축하여, 건물·도시·제품·의료 분야에 적용합니다.
셋째, 연구실에서 현장으로입니다. 연구는 논문에 머물지 않아야 합니다. 저의 연구는 콘서트홀 3건의 설계, 전기차 양산 사운드, VR 디지털 치료제 임상 검증에 실제 적용되었습니다.""")


# ═══════════════════════════════════════
# S6. RESEARCH OVERVIEW
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Research Overview — One Method, Four Impacts")
slide_num(s, 6)

pipeline = [("SENSE", "VR / AR\nAuralization"), ("UNDERSTAND", "Psychoacoustic\nEvaluation"),
            ("PREDICT", "AI / ML\nModeling"), ("APPLY", "Design &\nProduct")]
for i, (step, desc) in enumerate(pipeline):
    x = Inches(0.6) + Inches(i * 3.2)
    box = add_rounded_shape(s, x, y:=Inches(1.3), Inches(2.6), Inches(1.4), SAMSUNG_DARK)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = step; p.font.size = Pt(17); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = SAMSUNG_CYAN; p2.alignment = PP_ALIGN.CENTER
    if i < 3:
        tb(s, x + Inches(2.65), Inches(1.7), Inches(0.5), Inches(0.4),
           "→", sz=20, color=SAMSUNG_BLUE, bold=True, align=PP_ALIGN.CENTER)

impacts = [
    ("BUILDINGS", "Concert halls, apartments,\noffice environments", "12 SCI papers"),
    ("CITIES", "Urban parks, streets,\nsoundscape design", "7 SCI papers"),
    ("PRODUCTS", "EV AVAS, vehicle NVH,\nmass production", "2 SCI + JASA"),
    ("HEALTH", "Pneumonia AI diagnosis,\ndigital therapeutics", "3 SCI + TT"),
]
for i, (area, desc, output) in enumerate(impacts):
    x = Inches(0.6) + Inches(i * 3.2)
    add_rounded_shape(s, x, Inches(3.3), Inches(2.6), Inches(3.3), LIGHT_CARD)
    tb(s, x + Inches(0.2), Inches(3.5), Inches(2.2), Inches(0.4),
       area, sz=15, color=SAMSUNG_BLUE, bold=True, align=PP_ALIGN.CENTER)
    accent_bar(s, Inches(3.9), x + Inches(0.5), Inches(1.6), SAMSUNG_BLUE)
    tb(s, x + Inches(0.15), Inches(4.1), Inches(2.3), Inches(1.2),
       desc, sz=12, color=MID_GRAY, align=PP_ALIGN.CENTER, lsp=1.3)
    tb(s, x + Inches(0.15), Inches(5.5), Inches(2.3), Inches(0.4),
       output, sz=11, color=SAMSUNG_DARK, bold=True, align=PP_ALIGN.CENTER)

add_notes(s,
"""My research follows a four-stage pipeline: Sense, Understand, Predict, and Apply.
In the Sense stage, I use VR and AR auralization to recreate building acoustic environments in the laboratory.
In the Understand stage, I apply psychoacoustic evaluation methods with bio-signals to quantify how people perceive these environments.
In the Predict stage, I use AI and machine learning to build predictive models for environmental quality.
In the Apply stage, I translate research findings into actual building designs, urban planning guidelines, and commercial products.
This single methodology creates impact across four areas: Buildings, Cities, Products, and Health. Let me walk you through each one.""",
"""저의 연구는 네 단계의 파이프라인을 따릅니다: 감지, 이해, 예측, 적용입니다.
감지 단계에서는 VR/AR 오럴라이제이션으로 건축 음환경을 실험실에서 재현합니다.
이해 단계에서는 심리음향 평가와 생체신호로 인간의 환경 지각을 정량화합니다.
예측 단계에서는 AI/머신러닝으로 환경 품질 예측 모델을 구축합니다.
적용 단계에서는 연구 결과를 실제 건물 설계, 도시 계획, 상용 제품에 반영합니다.
이 하나의 방법론이 건물, 도시, 제품, 건강의 네 영역에 걸쳐 impact를 만듭니다. 각각을 설명드리겠습니다.""")


# ═══════════════════════════════════════
# S7. IMPACT 1: BUILDINGS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Impact 1: Buildings — Comfortable Built Environments")
slide_num(s, 7)

tb(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "Architectural Acoustics Design", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55))

items = [
    "•  3 Concert Halls: IFEZ Art Center (1,800 seats),",
    "   Bucheon Culture Arts Center (1,450), Pyeongtaek (1,315)",
    "•  Sound diffusion optimization with scale models + simulations [1]",
    "",
    "•  Floor impact noise: VR evaluation → criterion lowered by 6-7 dB [2]",
    "•  Open-plan office: Comfort-Content balance model [3][4]",
    "•  Water supply/drainage noise: VR annoyance assessment [5]",
]
tf = tb(s, Inches(0.8), Inches(1.75), Inches(5.8), Inches(3.0),
        items[0], sz=13, color=DARK_TEXT, lsp=1.3)
for it in items[1:]:
    c = DARK_TEXT
    ap(tf, it, sz=13, color=c, sa=2)

# Key message box
add_rounded_shape(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(4.5), LIGHT_CARD)
tb(s, Inches(7.3), Inches(1.4), Inches(5.0), Inches(0.4),
   "Key Contribution", sz=15, color=SAMSUNG_DARK, bold=True)
tb(s, Inches(7.3), Inches(1.9), Inches(5.0), Inches(3.5),
   "Established human perception-based\nevaluation criteria for building acoustics,\nreplacing conventional physical metrics.\n\n"
   "VR technology enables reproducible\nlaboratory experiments that closely\nmatch real-world acoustic experiences.\n\n"
   "Results directly informed Korean\nbuilding noise standards and\narchitectural design practice.",
   sz=13, color=MID_GRAY, lsp=1.4)

footnote(s,
    "[1] Jo & Jeon (2022) \"Optimizing sound diffusion in a concert hall using scale-model measurements and simulations\" JOBE 50, Q1\n"
    "[2] Jo & Jeon (2019) \"Downstairs resident classification characteristics for upstairs walking vibration noise...under VR\" B&E 150, Q1, Cited 50\n"
    "[3] Jeon, Jo* et al. (2022) \"Crossed effects of audio-visual environment on indoor soundscape perception for pleasant OPO\" B&E 207, Q1\n"
    "[4] Jo & Jeon (2022) \"Influence of indoor soundscape perception...on work-related quality with preference and productivity in OPO\" B&E 208, Q1\n"
    "[5] Jeon, Jo* et al. (2019) \"Subjective and objective evaluation of water-supply and drainage noises...using a head-mounted display\" Appl. Acoust. 148, Q1",
    y=Inches(6.35))

add_notes(s,
"""In the Buildings domain, I have two main contributions.
First, architectural acoustics design. I participated in the acoustical design of three concert halls totaling over 4,500 seats. I developed methods to optimize sound diffusion using both scale models and computer simulations, published in the Journal of Building Engineering.
Second, indoor sound environment evaluation. For floor impact noise in apartments, I used VR to conduct perceptual evaluations and discovered that the criterion for satisfaction should be 6 to 7 decibels stricter than existing Korean standards. This was published in Building and Environment with 50 citations.
For open-plan offices, I developed a Comfort-Content balance model showing the trade-off between acoustic preference and work productivity. And for water supply and drainage noise, I conducted VR-based annoyance assessments.
The key contribution is establishing human perception-based evaluation criteria that replace conventional physical metrics.""",
"""건물 영역에서의 주요 기여를 말씀드리겠습니다.
첫째, 건축음향 설계입니다. 총 4,565석 규모의 콘서트홀 3건의 음향설계에 참여하였고, 축소모형과 시뮬레이션을 활용한 음확산 최적화 방법론을 개발하여 Journal of Building Engineering에 게재하였습니다.
둘째, 실내 음환경 평가입니다. 아파트 층간소음에 대해 VR 기반 지각 평가를 수행한 결과, 기존 국내 기준 대비 6-7dB 엄격한 만족 기준이 필요함을 밝혔습니다. Building and Environment에 게재되어 50회 인용되었습니다.
오픈플랜 오피스에서는 음향 쾌적성과 업무 생산성 간의 trade-off를 발견한 Comfort-Content 밸런스 모델을 개발하였습니다.
핵심 기여는 기존 물리량 중심의 평가 기준을 인간 지각 기반 평가 기준으로 전환한 것입니다.""")


# ═══════════════════════════════════════
# S8. IMPACT 2: CITIES
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Impact 2: Cities — Pleasant Urban Sound Environments")
slide_num(s, 8)

tb(s, Inches(0.8), Inches(1.2), Inches(6.0), Inches(0.4),
   "Soundscape Assessment & Design", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55))

city = [
    "•  ISO 12913: Quantitative + qualitative protocol validation [1]",
    "•  Audio-visual interaction: Visual 90% contribution,",
    "   yet sound determines overall satisfaction [2] (Cited 240)",
    "•  Soundscape design index: Revised GSI/RSI/GLI/RLI",
    "   SEM model — \"Secure pleasantness, control eventfulness\" [3]",
    "•  Urban behavior: Human presence → dynamic soundscape [4]",
]
tf = tb(s, Inches(0.8), Inches(1.75), Inches(6.0), Inches(3.2),
        city[0], sz=13, color=DARK_TEXT, lsp=1.3)
for it in city[1:]:
    ap(tf, it, sz=13, color=DARK_TEXT, sa=3)

# Collaborations
tb(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.4),
   "International Collaboration", sz=17, color=SAMSUNG_DARK, bold=True)

collabs = [
    ("UCL, London", "Prof. Jian Kang", "Soundscape descriptors — 18-language\nvalidation (SATP Project) [5]"),
    ("Sorbonne, Paris", "Prof. J-D. Polack", "Paris-Seoul urban soundscape\ncomparative study (STAR Project)"),
    ("RWTH Aachen", "Prof. M. Vorlaender", "VR spatial audio rendering\n& ecological validity"),
]
for i, (inst, prof, desc) in enumerate(collabs):
    y = Inches(1.8) + Inches(i * 1.55)
    add_rounded_shape(s, Inches(7.2), y, Inches(5.3), Inches(1.3), LIGHT_CARD)
    tb(s, Inches(7.5), y + Inches(0.1), Inches(2.3), Inches(0.35),
       inst, sz=13, color=SAMSUNG_DARK, bold=True)
    tb(s, Inches(9.9), y + Inches(0.1), Inches(2.4), Inches(0.35),
       prof, sz=11, color=MID_GRAY)
    tb(s, Inches(7.5), y + Inches(0.45), Inches(4.8), Inches(0.7),
       desc, sz=11, color=MID_GRAY, lsp=1.25)

footnote(s,
    "[1] Jo & Jeon (2021) \"Compatibility of quantitative and qualitative data-collection protocols for urban soundscape evaluation\" SCS 74, Q1, IF 11.7\n"
    "[2] Jeon & Jo* (2020) \"Effects of audio-visual interactions on soundscape and landscape perception...satisfaction with the urban environment\" B&E 169, Q1, Cited 240\n"
    "[3] Jo & Jeon (2021) \"Overall environmental assessment in urban parks: Modelling audio-visual interaction with a SEM based on soundscape and landscape indices\" B&E 204, Q1, Cited 103\n"
    "[4] Jo & Jeon (2020) \"The influence of human behavioral characteristics on soundscape perception in urban parks\" L&UP 203, Q1, Cited 111\n"
    "[5] Aletta†, Jo et al. (2024) \"Soundscape descriptors in eighteen languages: Translation and validation through listening experiments\" Appl. Acoust. 224, Q1",
    y=Inches(6.35))

add_notes(s,
"""In the Cities domain, my main contributions are in soundscape assessment methodology and design guidelines.
For assessment, I validated the compatibility of quantitative and qualitative data collection protocols based on ISO 12913, published in Sustainable Cities and Society with impact factor 11.7.
A key finding was about audio-visual interaction. While visual information contributes 90% to environmental perception, it is actually the sound environment that determines overall satisfaction. This paper in Building and Environment has been cited 240 times, making it one of the most cited papers in the soundscape field.
I also developed revised soundscape and landscape indices with a structural equation model, establishing the design principle of "secure pleasantness while controlling eventfulness."
On the right, you can see my international collaborations. I worked with UCL on the 18-language validation of soundscape descriptors, with Sorbonne on Paris-Seoul comparative studies, and with RWTH Aachen on VR spatial audio rendering.""",
"""도시 영역에서의 기여를 말씀드립니다.
사운드스케이프 평가 방법론 측면에서, ISO 12913 기반의 정량·정성 데이터 수집 프로토콜의 호환성을 검증하였습니다.
핵심 발견은 시청각 상호작용에 관한 것입니다. 시각 정보가 환경 인지의 90%를 차지하지만, 전체 만족도를 결정하는 것은 실제로 음환경이었습니다. 이 논문은 Building and Environment에 게재되어 240회 인용되었으며, 사운드스케이프 분야에서 가장 많이 인용된 논문 중 하나입니다.
또한 수정된 사운드스케이프·경관 인덱스와 구조방정식 모델을 개발하여, "쾌적성을 확보하고 활력성을 제어한다"는 설계 원칙을 제시하였습니다.
오른쪽은 국제공동연구 현황입니다. UCL과는 18개국 사운드스케이프 기술어 검증, 소르본과는 파리-서울 비교연구, RWTH 아헨과는 VR 입체음향 연구를 수행하였습니다.""")


# ═══════════════════════════════════════
# S9. IMPACT 3: PRODUCTS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Impact 3: Products — From Research to Mass Production")
slide_num(s, 9)

tb(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "Hyundai Motor Company (4 Years)", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55))

prod = [
    "•  AVAS Brand Sound 2.0 Design & Tuning",
    "   Hyundai / Kia / Genesis (EV3, IONIQ 5)",
    "•  Soundscape-based competitor AVAS evaluation [1]",
    "   VR auralization with ambisonic reproduction system",
    "•  Sound environment quantification in EV cabins",
    "   Architectural acoustic theory → TDP methodology",
    "•  Vehicle NVH test: B-seg EV/ICE noise reduction (K3, EV3, IONIQ5)",
    "•  Industry-academia joint research PI",
    "   Seoul Nat'l Univ. + Chungnam Nat'l Univ. (NRF funded)",
]
tf = tb(s, Inches(0.8), Inches(1.75), Inches(6.0), Inches(3.8),
        prod[0], sz=13, color=DARK_TEXT, lsp=1.25)
for it in prod[1:]:
    ap(tf, it, sz=13, color=DARK_TEXT, sa=2)

# Pipeline
tb(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.4),
   "Research-to-Product Pipeline", sz=17, color=SAMSUNG_DARK, bold=True)

stages = [("Research", "Perception model &\nevaluation framework"),
          ("Evaluation", "Subjective listening test &\npsychoacoustic analysis"),
          ("Validation", "System-level verification &\nregulatory compliance"),
          ("Production", "Mass production tuning &\nquality management")]
for i, (stage, desc) in enumerate(stages):
    y = Inches(1.8) + Inches(i * 1.15)
    box = add_rounded_shape(s, Inches(7.2), y, Inches(1.8), Inches(0.8), SAMSUNG_DARK)
    tf_l = box.text_frame; tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_l.paragraphs[0].text = stage; tf_l.paragraphs[0].font.size = Pt(12)
    tf_l.paragraphs[0].font.color.rgb = WHITE; tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(s, Inches(9.2), y + Inches(0.05), Inches(3.3), Inches(0.7),
       desc, sz=11, color=MID_GRAY, lsp=1.25)
    if i < 3:
        tb(s, Inches(7.9), y + Inches(0.8), Inches(0.4), Inches(0.3),
           "↓", sz=13, color=SAMSUNG_BLUE, align=PP_ALIGN.CENTER)

# bottom note
tb(s, Inches(7.2), Inches(6.1), Inches(5.3), Inches(0.4),
   "\"Research → Evaluation → Validation → Production — the full cycle.\"",
   sz=13, color=SAMSUNG_DARK, bold=True)

footnote(s,
    "[1] Jeon, Jo & Hong (2026) \"Perceptual impact of electric vehicle acoustic alerting systems on urban soundscapes: using the ISO 12913 framework\" JASA — under review\n"
    "    Jo, Jeon & Hong (2025) \"Perceptual and satisfaction models for AVAS based on soundscape design concept\" HMG Summit Conference — Special Award",
    y=Inches(6.55))

add_notes(s,
"""In the Products domain, I bring four years of hands-on R&D experience at Hyundai Motor Company.
My primary role was designing AVAS — the Acoustic Vehicle Alerting System — for electric vehicles. I led the Brand Sound 2.0 design for Hyundai, Kia, and Genesis models including EV3 and IONIQ 5.
I also developed a soundscape-based framework for evaluating competitor AVAS sounds using VR auralization, and proposed a sound environment quantification methodology for EV cabins based on architectural acoustic theory.
On the NVH side, I worked on noise reduction for B-segment vehicles and managed industry-academia joint research with Seoul National University and Chungnam National University.
On the right, you can see the research-to-product pipeline I experienced: from perception model development, through listening tests and system verification, all the way to mass production tuning. This end-to-end experience is what distinguishes my research approach.
Our AVAS soundscape research was submitted to the Journal of the Acoustical Society of America and is currently under review.""",
"""제품 영역에서는 현대자동차 연구소 4년간의 실무 R&D 경험을 말씀드립니다.
주요 업무는 전기차 보행자 경고음(AVAS) 설계였습니다. 현대·기아·제네시스의 EV3, 아이오닉5 등 브랜드 사운드 2.0 디자인을 주도하였습니다.
또한 VR 오럴라이제이션을 활용한 사운드스케이프 기반 경쟁사 AVAS 평가 프레임워크를 개발하였고, 건축음향 이론을 적용한 전기차 캐빈 음환경 정량화 방법론을 제안하였습니다.
NVH 측면에서는 준중형 전기차/내연기관 차량의 소음저감 연구를 수행하였고, 서울대학교·충남대학교와의 산학공동연구를 기업 측 PI로서 운영하였습니다.
오른쪽은 제가 경험한 연구→제품 파이프라인입니다. 지각 모델 개발부터 청감 평가, 시스템 검증, 양산 튜닝까지 전 과정을 경험하였습니다.
AVAS 사운드스케이프 연구는 JASA에 투고하여 현재 심사 중입니다.""")


# ═══════════════════════════════════════
# S10. IMPACT 4: HEALTH
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Impact 4: Health — Sound for Human Well-being")
slide_num(s, 10)

tb(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "AI-based Respiratory Disease Diagnosis", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55))

ai_items = [
    "•  Cough sound classification using LSTM / CNN",
    "   → Pneumonia diagnosis accuracy: 84.9% [1]",
    "•  RIR-based data augmentation for deep learning",
    "   → Patent (Korea + US No. 18273592) [2]",
    "   → Technology transfer: $50,000 to Hanyang S&A",
]
tf = tb(s, Inches(0.8), Inches(1.75), Inches(5.8), Inches(2.0),
        ai_items[0], sz=13, color=DARK_TEXT, lsp=1.3)
for it in ai_items[1:]:
    c = SAMSUNG_BLUE if "→" in it else DARK_TEXT
    ap(tf, it, sz=13, color=c, sa=2)

tb(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.4),
   "VR Soundscape Digital Therapeutics", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(3.95))

vr_items = [
    "•  VR sound therapy for mental illness",
    "   (Depression, schizophrenia — HYU Hospital)",
    "•  Psycho-physiological restoration quantified",
    "   via EEG, HRV, and eye-tracking [3][4]",
]
tf2 = tb(s, Inches(0.8), Inches(4.15), Inches(5.8), Inches(1.8),
         vr_items[0], sz=13, color=DARK_TEXT, lsp=1.3)
for it in vr_items[1:]:
    ap(tf2, it, sz=13, color=DARK_TEXT, sa=2)

# Right message
add_rounded_shape(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(4.5), LIGHT_CARD)
tb(s, Inches(7.3), Inches(1.4), Inches(5.0), Inches(0.4),
   "Expanding Building Acoustics", sz=15, color=SAMSUNG_DARK, bold=True)
tb(s, Inches(7.3), Inches(1.9), Inches(5.0), Inches(3.5),
   "Building acoustics methodologies\n— sound analysis, spatial audio,\npsychoacoustic evaluation —\n\n"
   "extend to:\n\n"
   "✓  Healthcare (disease diagnosis)\n"
   "✓  Mental health (digital therapeutics)\n"
   "✓  Wellness (restorative environments)\n\n"
   "Demonstrates broad applicability of\nperception-driven sound research\nbeyond traditional building engineering.",
   sz=13, color=MID_GRAY, lsp=1.35)

footnote(s,
    "[1] Chung, Jo et al. (2021) \"Diagnosis of pneumonia by cough sounds analyzed with statistical features and AI\" Sensors 21(21), SCI-E\n"
    "[2] Jeon, Jo et al. (2021/2024) Patent: \"Improvement of deep learning based sound classifier performance with data augmentation with RIR\" KR+US 18273592, TT ₩50M\n"
    "[3] Jeon, Jo* & Lee (2023) \"Psycho-physiological restoration with audio-visual interactions through VR simulations of soundscape and landscape\" SCS 99, Q1, Cited 74\n"
    "[4] Jo et al. (2022) \"Effect of noise sensitivity on psychophysiological response through monoscopic 360 video and stereoscopic sound\" Scientific Report 12, Q1",
    y=Inches(6.35))

add_notes(s,
"""The fourth impact area is Health, where building acoustics methodologies extend to healthcare applications.
In AI-based diagnosis, I contributed to a study using LSTM and CNN to classify cough sounds for pneumonia diagnosis, achieving 84.9% accuracy. To address data scarcity, I developed a room impulse response-based data augmentation technique that improved deep learning classifier performance. This technology was patented both in Korea and the US, and transferred to industry for 50,000 dollars.
In digital therapeutics, I worked on VR-based soundscape therapy for mental illness patients including depression and schizophrenia, in collaboration with Hanyang University Hospital. We quantified psycho-physiological restoration effects using EEG, heart rate variability, and eye-tracking.
This demonstrates that building acoustics methodologies have broad applicability beyond traditional engineering — into healthcare, mental health, and wellness.""",
"""네 번째 영역인 건강 분야에서는 건축음향 방법론의 의료 확장을 말씀드립니다.
AI 기반 진단에서는 LSTM과 CNN을 활용한 기침소리 분류로 84.9%의 폐렴 진단 정확도를 달성하였습니다. 데이터 부족 문제를 해결하기 위해 공간 임펄스 응답 기반 데이터 증강 기술을 개발하였고, 이 기술은 한국과 미국에서 특허 등록되었으며 5천만원에 기술이전되었습니다.
디지털 치료제에서는 한양대학교 병원과 협력하여 우울증, 조현병 환자 대상 VR 사운드스케이프 치료 콘텐츠를 개발하고, EEG, 심박변이도, 시선추적으로 심리생리학적 회복 효과를 정량화하였습니다.
이는 건축음향의 방법론이 전통적 공학을 넘어 의료, 정신건강, 웰니스까지 확장될 수 있음을 보여줍니다.""")


# ═══════════════════════════════════════
# S11. GLOBAL NETWORK
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Global Network & Scholarly Activities")
slide_num(s, 11)

tb(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "International Joint Research Projects", sz=17, color=SAMSUNG_DARK, bold=True)
joints = [
    ("SATP Project (2020-2022)", "Soundscape Attributes Translation — 18 languages\nUCL, Stockholm, TU Berlin, McGill, NUS, etc."),
    ("STAR Project (2018-2020)", "France-Korea urban soundscape comparison\nSorbonne University, Paris"),
    ("CHIC Project (2016-2019)", "VR 3D audio environment modeling\nRWTH Aachen University, KIST"),
]
for i, (t, d) in enumerate(joints):
    y = Inches(1.7) + Inches(i * 1.45)
    add_rounded_shape(s, Inches(0.8), y, Inches(5.5), Inches(1.2), LIGHT_CARD)
    tb(s, Inches(1.1), y + Inches(0.1), Inches(5.0), Inches(0.3), t, sz=13, color=SAMSUNG_DARK, bold=True)
    tb(s, Inches(1.1), y + Inches(0.4), Inches(5.0), Inches(0.7), d, sz=11, color=MID_GRAY, lsp=1.25)

tb(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.4),
   "Awards & Scholarly Activities", sz=17, color=SAMSUNG_DARK, bold=True)

awards = [
    "Awards (selected):",
    "  •  EAA Best Paper & Presentation — ICA 2019, Aachen",
    "  •  I-INCE Young Professional Award — Inter-Noise 2020",
    "  •  HMG Special Award for Research Paper — 2025",
    "  •  Promising Scientist Award — KSNVE 2022",
    "  •  Ph.D. Excellence Paper Award — Hanyang 2022",
    "",
    "Journal Reviewer: 17 international journals",
    "  B&E, L&UP, SCS, Applied Acoustics, IEEE Access, etc.",
    "",
    "Keynote: Urban Sound Symposium, E-congress 2021",
    "",
    "Affiliations: ASA | AES | KSNVE | ASK | AIK | SAREK",
]
tf = tb(s, Inches(7.2), Inches(1.7), Inches(5.3), Inches(4.8),
        awards[0], sz=13, color=SAMSUNG_DARK, bold=True, lsp=1.2)
for it in awards[1:]:
    b = bool(it and not it.startswith("  ") and not it.startswith("•"))
    ap(tf, it, sz=12 if it.startswith("  ") else 13, color=DARK_TEXT, bold=b, sa=1)

add_notes(s,
"""This slide summarizes my global network and scholarly activities.
I have been involved in three major international joint research projects: the SATP project with UCL on 18-language soundscape descriptor validation, the STAR project with Sorbonne on Paris-Seoul urban soundscape comparison, and the CHIC project with RWTH Aachen on VR 3D audio modeling.
In terms of awards, I received the EAA Best Paper and Presentation Award at the International Congress on Acoustics in 2019, the I-INCE Young Professional Award in 2020, and most recently the Hyundai Motor Group Special Award for research in 2025.
I serve as a reviewer for 17 international journals and delivered a keynote at the Urban Sound Symposium in 2021. I am also a member of major acoustical societies in both Korea and internationally.""",
"""국제 네트워크와 학술 활동을 요약합니다.
3건의 주요 국제공동연구를 수행하였습니다. UCL과의 SATP 프로젝트(18개국 사운드스케이프 기술어 검증), 소르본과의 STAR 프로젝트(파리-서울 비교연구), RWTH 아헨과의 CHIC 프로젝트(VR 3D 오디오 모델링)입니다.
수상으로는 2019년 ICA 학술대회 EAA 최우수 논문·발표상, 2020년 I-INCE 유망연구자상, 2025년 현대자동차그룹 학술대회 특별상 등이 있습니다.
17개 국제 저널의 리뷰어로 활동하고 있으며, 2021년 Urban Sound Symposium에서 기조강연을 하였습니다.""")


# ═══════════════════════════════════════
# S12. WHAT INDUSTRY TAUGHT ME
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "What Industry Taught Me")
slide_num(s, 12)

lessons = [
    ("Research-to-Product\nPipeline",
     "Experienced the full cycle:\nperception model → listening test →\nsystem verification → mass production.\n\n"
     "Research validated not only through\npeer review, but through real-world\nproduct deployment."),
    ("Industry-Ready\nResearch Network",
     "Managed industry-academia joint\nresearch as PI from the corporate side.\n\n"
     "Ready to initiate joint research with\nHyundai Motor, Samsung, and\nconstruction companies from Day 1."),
    ("Real-World Problem\nDefinition",
     "Industry problems are the best\nresearch questions.\n\n"
     "\"How should AVAS sound in future\ncities?\" — from the factory floor\nto a JASA paper."),
]
for i, (title, desc) in enumerate(lessons):
    x = Inches(0.6) + Inches(i * 4.2)
    add_rounded_shape(s, x, Inches(1.3), Inches(3.9), Inches(4.8), LIGHT_CARD)
    num_box = add_rounded_shape(s, x + Inches(0.2), Inches(1.5), Inches(0.5), Inches(0.5), SAMSUNG_DARK)
    tf_n = num_box.text_frame; tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.paragraphs[0].text = str(i+1); tf_n.paragraphs[0].font.size = Pt(18)
    tf_n.paragraphs[0].font.color.rgb = WHITE; tf_n.paragraphs[0].font.bold = True
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(s, x + Inches(0.85), Inches(1.5), Inches(2.8), Inches(0.6),
       title, sz=15, color=SAMSUNG_DARK, bold=True, lsp=1.1)
    tb(s, x + Inches(0.3), Inches(2.3), Inches(3.3), Inches(3.5),
       desc, sz=13, color=MID_GRAY, lsp=1.4)

tb(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
   "\"This experience shaped how I define, validate, and deliver research.\"",
   sz=16, color=SAMSUNG_DARK, bold=True, align=PP_ALIGN.CENTER)

add_notes(s,
"""Before moving to my future plans, let me share three key lessons from my four years in industry.
First, the Research-to-Product Pipeline. I experienced the entire cycle from building a perception model, through listening tests and system verification, to mass production tuning. My research was validated not just through peer review but through actual product deployment in commercial vehicles.
Second, an Industry-Ready Research Network. I managed joint research projects from the corporate side as PI. This means I can initiate industry-academia collaborations from Day 1, with partners like Hyundai Motor and Samsung already in my network.
Third, Real-World Problem Definition. The best research questions come from industry. The question "how should AVAS sound in future cities?" came from the factory floor and became a paper submitted to JASA.
This experience fundamentally shaped how I define, validate, and deliver research.""",
"""향후 계획으로 넘어가기 전에, 산업체 4년에서 배운 세 가지 교훈을 말씀드립니다.
첫째, 연구에서 제품까지의 파이프라인입니다. 지각 모델 구축부터 청감 평가, 시스템 검증, 양산 튜닝까지 전 과정을 경험하였습니다.
둘째, 산업체 연구 네트워크입니다. 기업 측 PI로서 산학과제를 운영한 경험이 있어, 부임 첫날부터 현대차, 삼성 등과 산학협력을 시작할 수 있습니다.
셋째, 현장 기반 연구 문제 정의입니다. "미래 도시에서 AVAS는 어떤 소리여야 하는가?" — 현장에서 나온 질문이 JASA 논문이 되었습니다.
이 경험이 제가 연구를 정의하고, 검증하고, 전달하는 방식을 근본적으로 바꾸어 놓았습니다.""")


# ═══════════════════════════════════════
# S13. SENSE Lab Vision
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SAMSUNG_DARK)
slide_num(s, 13)

tb(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
   "FUTURE RESEARCH PLAN", sz=14, color=SAMSUNG_CYAN, bold=True)
tb(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
   "SENSE Lab", sz=52, color=WHITE, bold=True)
tb(s, Inches(0.8), Inches(2.6), Inches(11), Inches(0.5),
   "Sound Environment aNd Sensory Engineering Laboratory",
   sz=20, color=SAMSUNG_CYAN)
tb(s, Inches(0.8), Inches(3.6), Inches(11), Inches(0.5),
   "Completing the building environment research at Hanyang:",
   sz=17, color=WHITE)

pillars_data = [
    ("Thermal\nEnvironment", "(Existing)", RGBColor(0x15, 0x30, 0x55)),
    ("Air Quality\nEnvironment", "(Existing)", RGBColor(0x15, 0x30, 0x55)),
    ("Sound\nEnvironment", "(SENSE Lab)", SAMSUNG_BLUE),
]
for i, (name, status, clr) in enumerate(pillars_data):
    x = Inches(0.8) + Inches(i * 3.5)
    box = add_rounded_shape(s, x, Inches(4.3), Inches(3.0), Inches(1.2), clr)
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(15); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = status; p2.font.size = Pt(11)
    p2.font.color.rgb = SAMSUNG_CYAN; p2.alignment = PP_ALIGN.CENTER

tb(s, Inches(4.5), Inches(5.55), Inches(4.3), Inches(0.35),
   "▼              ▼              ▼", sz=13, color=SAMSUNG_CYAN, align=PP_ALIGN.CENTER)

ibox = add_rounded_shape(s, Inches(2.5), Inches(5.9), Inches(8.3), Inches(0.85), SAMSUNG_BLUE)
tf_i = ibox.text_frame; tf_i.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_i.paragraphs[0].text = "Intelligent Building Environment System (IBES)"
tf_i.paragraphs[0].font.size = Pt(18); tf_i.paragraphs[0].font.color.rgb = WHITE
tf_i.paragraphs[0].font.bold = True; tf_i.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(s,
"""Now I present my future research vision: SENSE Lab — Sound Environment and Sensory Engineering Laboratory.
Currently, Hanyang University's architectural engineering department has established research groups in thermal environment and air quality. However, the sound environment — one of the three fundamental pillars of building environment — remains vacant.
SENSE Lab will complete this triangle. By adding sound to thermal and air quality, we can build a truly Intelligent Building Environment System — or IBES — that integrates all three domains through AI and data-driven approaches.
This is not just about adding another research area. It's about enabling cross-domain synergy that makes the whole department's building environment research greater than the sum of its parts.""",
"""이제 향후 연구 비전을 말씀드립니다. SENSE Lab — Sound Environment aNd Sensory Engineering Laboratory입니다.
현재 한양대학교 건축공학부에는 열환경과 공기질 분야의 연구 그룹이 있습니다. 그러나 건축환경의 3대 축 중 하나인 음환경은 공백 상태입니다.
SENSE Lab이 이 삼각형을 완성합니다. 열과 공기에 소리를 더하여, AI와 데이터 기반으로 세 영역을 통합하는 진정한 지능형 건축환경 시스템(IBES)을 구축할 수 있습니다.
이것은 단순히 연구 분야를 추가하는 것이 아니라, 학과 전체의 건축환경 연구가 시너지를 만들 수 있는 기반을 마련하는 것입니다.""")


# ═══════════════════════════════════════
# S14. PHASE 1+2
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Research Plan — Phase 1 & Phase 2")
slide_num(s, 14)

for phase, title, items, col_x in [
    ("PHASE 1  |  Year 1–3", "AI-driven Building Acoustic\nEnvironment Diagnosis",
     ["① IoT acoustic sensor network for indoor/outdoor\n    sound environment monitoring",
      "② ML-based sound environment auto-classification\n    and anomaly detection",
      "③ VR simulation platform for design-stage\n    acoustic previewing",
      "", "Target Funding:",
      "  •  NRF Young/Mid-career (₩200–300M/yr)",
      "  •  MSIT Smart City Program",
      "  •  Hyundai Motor joint research",
      "", "Output: 3–4 SCI/yr, 3–4 graduate students"], Inches(0.6)),
    ("PHASE 2  |  Year 3–5", "Mobility-Urban Sound\nEnvironment Interaction",
     ["① Long-term urban sound environment monitoring\n    in the EV transition era",
      "② AVAS + urban soundscape integrated design\n    framework and guidelines",
      "③ Digital twin-based sound environment\n    prediction and simulation",
      "", "Target Funding:",
      "  •  NRF Leading Research Center (₩500–1B/yr)",
      "  •  MOLIT / MOE urban noise policy",
      "  •  Samsung — Spatial Audio R&D",
      "", "Output: ISO standard, 6–8 students, post-doc"], Inches(6.8)),
]:
    add_rounded_shape(s, col_x, Inches(1.15), Inches(6.0), Inches(5.8), LIGHT_CARD)
    tb(s, col_x + Inches(0.2), Inches(1.25), Inches(3), Inches(0.35),
       phase, sz=14, color=SAMSUNG_BLUE, bold=True)
    tb(s, col_x + Inches(0.2), Inches(1.65), Inches(5.5), Inches(0.6),
       title, sz=18, color=SAMSUNG_DARK, bold=True, lsp=1.1)
    tf = tb(s, col_x + Inches(0.2), Inches(2.5), Inches(5.5), Inches(4.0),
            items[0], sz=12, color=DARK_TEXT, lsp=1.25)
    for it in items[1:]:
        b = bool(it.startswith("Target") or it.startswith("Output"))
        c = SAMSUNG_DARK if b else (SAMSUNG_BLUE if it.startswith("  •") else DARK_TEXT)
        ap(tf, it, sz=12, color=c, bold=b, sa=1)

add_notes(s,
"""My research plan consists of three phases. Let me present Phase 1 and Phase 2.
Phase 1, Years 1 through 3, focuses on building the foundation: an AI-driven building acoustic environment diagnosis system. This includes deploying IoT acoustic sensor networks, developing ML-based sound classification algorithms, and creating a VR simulation platform for previewing acoustics at the design stage. Target funding is 200 to 300 million won per year from NRF and industry partners.
Phase 2, Years 3 through 5, expands to Mobility-Urban Sound Environment Interaction. As electric vehicles transform city soundscapes, we need long-term monitoring systems and integrated AVAS-urban soundscape design frameworks. I also plan to develop digital twin-based prediction tools. Target funding scales to 500 million to 1 billion won per year, with partnerships including Samsung for spatial audio research.""",
"""연구 계획은 3단계로 구성됩니다. Phase 1과 2를 먼저 말씀드립니다.
Phase 1은 1-3년차로, AI 기반 건축 음환경 진단 시스템 구축에 집중합니다. IoT 음향 센서 네트워크 배치, ML 기반 음환경 자동 분류, VR 시뮬레이션 플랫폼 개발을 포함합니다. 목표 연구비는 연 2-3억원입니다.
Phase 2는 3-5년차로, 모빌리티-도시 음환경 상호작용 연구로 확장합니다. 전기차 시대의 도시 음환경 변화를 장기 모니터링하고, AVAS-도시 사운드스케이프 통합 설계 프레임워크를 개발합니다. 디지털 트윈 기반 예측 도구도 구축합니다. 목표 연구비는 연 5-10억원이며, 삼성과의 공간 오디오 공동연구를 포함합니다.""")


# ═══════════════════════════════════════
# S15. PHASE 3 + FUNDING
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Research Plan — Phase 3 & Funding Strategy")
slide_num(s, 15)

tb(s, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.35),
   "PHASE 3  |  Year 5+", sz=14, color=SAMSUNG_BLUE, bold=True)
tb(s, Inches(0.8), Inches(1.55), Inches(6), Inches(0.5),
   "Multi-modal Intelligent Building Env. Platform", sz=18, color=SAMSUNG_DARK, bold=True)

p3 = [
    "① Integrated sensing & control: Sound + Thermal + Air Quality",
    "   → Joint research with existing Hanyang AE faculty",
    "② Human-centered building environment optimization",
    "   → Wellness, productivity, energy efficiency simultaneously",
    "③ International research hub (UCL / RWTH / Sorbonne)",
    "   → ISO 12913 standard expansion and leadership",
]
tf = tb(s, Inches(0.8), Inches(2.1), Inches(6.2), Inches(2.5),
        p3[0], sz=13, color=DARK_TEXT, lsp=1.3)
for it in p3[1:]:
    c = SAMSUNG_BLUE if "→" in it else DARK_TEXT
    ap(tf, it, sz=13, color=c, sa=2)

# Funding roadmap
tb(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.4),
   "Funding Roadmap", sz=17, color=SAMSUNG_DARK, bold=True)

roadmap = [
    ("Year 1–2", "~₩400M/yr", "NRF Young Researcher + Hyundai joint"),
    ("Year 3–4", "~₩700M/yr", "NRF Mid-career + MOLIT + Industry ×3"),
    ("Year 5+", "~₩1B+/yr", "Leading Research Center + Int'l grants"),
]
for i, (yr, amt, desc) in enumerate(roadmap):
    y = Inches(1.8) + Inches(i * 1.2)
    add_rounded_shape(s, Inches(7.2), y, Inches(5.3), Inches(0.95), LIGHT_CARD)
    tb(s, Inches(7.5), y + Inches(0.08), Inches(1.5), Inches(0.35), yr, sz=13, color=SAMSUNG_DARK, bold=True)
    tb(s, Inches(9.2), y + Inches(0.08), Inches(1.5), Inches(0.35), amt, sz=13, color=SAMSUNG_BLUE, bold=True)
    tb(s, Inches(7.5), y + Inches(0.4), Inches(4.8), Inches(0.45), desc, sz=11, color=MID_GRAY)

tb(s, Inches(7.2), Inches(5.4), Inches(5.3), Inches(0.35),
   "Ready to Launch from Day 1:", sz=13, color=SAMSUNG_DARK, bold=True)
ready = [
    "•  Hyundai Motor — existing relationship (AVAS / NVH)",
    "•  Samsung — spatial audio & immersive media capability",
    "•  NRF — 11 funded projects on record",
    "•  UCL / RWTH / Sorbonne — active collaborations",
]
tf2 = tb(s, Inches(7.2), Inches(5.7), Inches(5.3), Inches(1.5),
         ready[0], sz=12, color=MID_GRAY, lsp=1.25)
for r in ready[1:]:
    ap(tf2, r, sz=12, color=MID_GRAY, sa=1)

add_notes(s,
"""Phase 3, from Year 5 onwards, aims for full integration. The goal is a multi-modal intelligent building environment platform that combines sound, thermal, and air quality sensing and control. This will be achieved through joint research with existing faculty members in the department.
I also plan to establish Hanyang as an international research hub leveraging my existing network, and take a leadership role in expanding ISO 12913 soundscape standards.
On the right is the funding roadmap. Starting at approximately 400 million won per year in Years 1-2, scaling to 700 million in Years 3-4, and targeting over 1 billion won annually from Year 5.
Critically, I can start from Day 1 with ready-to-launch partnerships: Hyundai Motor for AVAS and NVH research, Samsung for spatial audio, NRF where I have 11 funded projects on record, and active international collaborations with UCL, RWTH Aachen, and Sorbonne.""",
"""Phase 3은 5년차 이후로, 완전한 통합을 목표로 합니다. 음·열·공기질 센싱과 제어를 결합하는 멀티모달 지능형 건축환경 플랫폼을 구축하며, 학과 내 기존 교수님들과의 공동연구로 달성합니다.
또한 기존 네트워크를 활용하여 한양대를 국제 연구 허브로 자리매김하고, ISO 12913 사운드스케이프 표준 확장에서 리더십을 확보하겠습니다.
오른쪽은 과제 수주 로드맵입니다. 1-2년차 연 4억, 3-4년차 연 7억, 5년차 이후 연 10억 이상을 목표로 합니다.
핵심적으로, 부임 첫날부터 즉시 가동 가능한 파트너십이 있습니다: 현대차(AVAS/NVH), 삼성(공간음향), NRF(기존 11건 수행 이력), UCL/RWTH/Sorbonne(진행 중 공동연구)입니다.""")


# ═══════════════════════════════════════
# S16. FUNDING DETAIL
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Funding Strategy — Detailed Plan")
slide_num(s, 16)

tb(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "Government & Public Funding", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55))

gov = ["National Research Foundation (NRF):",
       "  •  Young/Mid-career: Soundscape AI (₩100–200M/yr)",
       "  •  International Joint: UCL, Sorbonne (₩120M/yr each)",
       "  •  Basic Research Lab: Building Env. Integration (₩500M/3yr)",
       "", "Government Agencies:",
       "  •  MSIT: Smart city sound monitoring",
       "  •  MOLIT / KAIA: Urban noise management policy",
       "  •  MOE: Environmental noise impact assessment",
       "  •  LH Corporation: Residential noise standards"]
tf = tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.0),
        gov[0], sz=13, color=SAMSUNG_DARK, bold=True, lsp=1.2)
for g in gov[1:]:
    b = bool(g and not g.startswith("  ") and g != "")
    ap(tf, g, sz=12 if g.startswith("  ") else 13, color=SAMSUNG_DARK if b else DARK_TEXT, bold=b, sa=2)

tb(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.4),
   "Industry Collaboration", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55), Inches(7.2))

inds = [
    ("Hyundai Motor Company", "AVAS sound design, vehicle cabin acoustics\n(Existing — joint research since 2023)"),
    ("Samsung Research", "Spatial audio quality, immersive media\n(Demonstrated capability in portfolio)"),
    ("Construction / Architecture", "Building acoustics consulting, noise assessment\n(Precedent: FURSYS ₩60M project)"),
]
for i, (co, desc) in enumerate(inds):
    y = Inches(1.8) + Inches(i * 1.55)
    add_rounded_shape(s, Inches(7.2), y, Inches(5.3), Inches(1.3), LIGHT_CARD)
    tb(s, Inches(7.5), y + Inches(0.1), Inches(4.8), Inches(0.35), co, sz=13, color=SAMSUNG_DARK, bold=True)
    tb(s, Inches(7.5), y + Inches(0.45), Inches(4.8), Inches(0.7), desc, sz=11, color=MID_GRAY, lsp=1.25)

add_notes(s,
"""This slide provides more detail on my funding strategy.
For government funding, I plan to target NRF programs at multiple levels: Young or Mid-career Researcher grants for soundscape AI research at 100-200 million won per year, International Joint Research with UCL and Sorbonne, and a Basic Research Lab grant for building environment integration.
I also plan to apply to government agencies including MSIT for smart city programs, MOLIT for urban noise policy, and LH Corporation for residential noise standards.
For industry collaboration, I have three main channels. Hyundai Motor Company, where I have an existing relationship and ongoing joint research since 2023. Samsung Research, where I have demonstrated spatial audio capability. And construction and architecture firms for building acoustics consulting, building on a precedent project with FURSYS worth 60 million won.""",
"""과제 수주 전략을 상세히 말씀드립니다.
정부 과제로는 NRF 다단계 프로그램을 타겟합니다: 신진/중견연구자 과제(사운드스케이프 AI, 연 1-2억), 국제공동연구(UCL/소르본, 각 연 1.2억), 기초연구실(건축환경 통합, 5억/3년)입니다.
정부 기관으로는 과기정통부 스마트시티, 국토부 도시소음 정책, LH 주거소음 기준 사업 등을 지원할 계획입니다.
산학협력은 세 채널입니다. 현대자동차(기존 관계, 2023년부터 공동연구 진행 중), 삼성리서치(공간음향 역량 보유), 건설/건축사무소(건축음향 컨설팅, 퍼시스 6천만원 과제 선례)입니다.""")


# ═══════════════════════════════════════
# S17. TEACHING PHILOSOPHY
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Teaching Philosophy")
slide_num(s, 17)

tb(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.7),
   "\"Educating future-ready architectural engineers\nwho understand both technology and human experience\"",
   sz=18, color=SAMSUNG_DARK, bold=True, align=PP_ALIGN.CENTER, lsp=1.3)

tp = [
    ("Theory → Practice",
     "Connect theory to real-world applications\nthrough industry case studies.\n\n"
     "Use Hyundai Motor and concert hall\nprojects as teaching materials.\n\n"
     "VR/AR simulation labs for immersive\nlearning experiences."),
    ("Data Literacy",
     "Equip architectural engineers with\nAI and data analysis capabilities.\n\n"
     "Python/MATLAB-based acoustic data\nprocessing and visualization.\n\n"
     "ML fundamentals for building\nperformance optimization."),
    ("Global Exposure",
     "International conference presentation\nrequired for all graduate students.\n\n"
     "Short-term visits to UCL, RWTH,\nSorbonne partner institutions.\n\n"
     "Bilingual research training for\nglobal careers."),
]
for i, (title, desc) in enumerate(tp):
    x = Inches(0.8) + Inches(i * 4.1)
    add_rounded_shape(s, x, Inches(2.2), Inches(3.8), Inches(4.5), LIGHT_CARD)
    tb(s, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.45),
       title, sz=16, color=SAMSUNG_DARK, bold=True)
    accent_bar(s, Inches(2.8), x + Inches(0.3), Inches(1.5), SAMSUNG_BLUE)
    tb(s, x + Inches(0.3), Inches(3.0), Inches(3.2), Inches(3.2),
       desc, sz=13, color=MID_GRAY, lsp=1.4)

add_notes(s,
"""My teaching philosophy centers on educating future-ready architectural engineers who understand both technology and human experience.
Three principles guide my approach.
First, Theory to Practice. I connect classroom theory to real-world applications using actual industry cases. Students will learn from Hyundai Motor projects and concert hall designs, and use VR simulation labs for immersive learning.
Second, Data Literacy. Modern architectural engineers need AI and data analysis skills. I will integrate Python and MATLAB-based acoustic data processing into the curriculum, along with machine learning fundamentals for building performance optimization.
Third, Global Exposure. All graduate students will be required to present at international conferences. I will facilitate short-term visits to our partner institutions — UCL, RWTH Aachen, and Sorbonne — and provide bilingual research training to prepare students for global careers.""",
"""저의 교육 철학은 기술과 인간 경험을 모두 이해하는 미래형 건축 엔지니어 양성입니다.
세 가지 원칙이 있습니다.
첫째, 이론에서 실천으로입니다. 현대자동차 프로젝트와 콘서트홀 설계 사례를 교재로 활용하고, VR 시뮬레이션 실습실에서 몰입형 학습 경험을 제공합니다.
둘째, 데이터 리터러시입니다. 현대 건축 엔지니어에게 AI와 데이터 분석 역량이 필수입니다. Python/MATLAB 기반 음향 데이터 처리와 머신러닝 기초를 교과과정에 통합합니다.
셋째, 글로벌 노출입니다. 대학원생 전원의 국제학회 발표를 의무화하고, UCL/RWTH/소르본 파트너 기관 단기 방문을 지원하며, 글로벌 커리어를 위한 이중언어 연구 훈련을 제공합니다.""")


# ═══════════════════════════════════════
# S18. TEACHING PLAN + MENTORING
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Teaching Plan & Student Mentoring")
slide_num(s, 18)

tb(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "Course Portfolio", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55))

tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.35),
   "Courses to Inherit:", sz=13, color=SAMSUNG_BLUE, bold=True)
tf = tb(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.2),
        "•  Building Environmental Engineering — UG", sz=12, color=DARK_TEXT, lsp=1.3)
ap(tf, "•  Architectural Acoustics — UG", sz=12, color=DARK_TEXT, sa=3)
ap(tf, "•  Probability & Statistics — UG", sz=12, color=DARK_TEXT, sa=3)

tb(s, Inches(0.8), Inches(3.3), Inches(5.5), Inches(0.35),
   "New Courses:", sz=13, color=SAMSUNG_BLUE, bold=True)
tf2 = tb(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(2.5),
         "•  Smart Building Environment Systems — Grad", sz=12, color=DARK_TEXT, lsp=1.2)
ap(tf2, "   IoT + ML + VR for building environment design", sz=11, color=MID_GRAY, sa=2)
ap(tf2, "•  AI for Built Environment — Grad", sz=12, color=DARK_TEXT, sa=2)
ap(tf2, "   Python/MATLAB acoustic data analysis lab", sz=11, color=MID_GRAY, sa=2)
ap(tf2, "•  Sustainable Building & Ecology — UG", sz=12, color=DARK_TEXT, sa=2)
ap(tf2, "   LEED/WELL standards, green building practice", sz=11, color=MID_GRAY, sa=2)

# Mentoring
tb(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.4),
   "Student Mentoring", sz=17, color=SAMSUNG_DARK, bold=True)
accent_bar(s, Inches(1.55), Inches(7.2))

for label, items, y_start in [
    ("Undergraduate:", ["•  Capstone design → industry project linkage",
                         "•  VR/BIM hands-on lab experiences",
                         "•  Early research exposure program"], Inches(1.8)),
    ("Graduate:", ["•  Target: 1 SCI paper + 1 patent per student",
                    "•  Int'l conference presentation required",
                    "•  Short-term visit to UCL / RWTH / Sorbonne",
                    "•  Weekly seminar + writing workshop"], Inches(3.3)),
    ("Career Paths:", ["•  Hyundai Motor / Kia — NVH",
                        "•  Samsung Research — Audio Lab",
                        "•  Construction & architecture firms",
                        "•  Government institutes (KICT, KRICT)"], Inches(5.1)),
]:
    tb(s, Inches(7.2), y_start, Inches(5.3), Inches(0.35), label, sz=13, color=SAMSUNG_BLUE, bold=True)
    tf_m = tb(s, Inches(7.2), y_start + Inches(0.3), Inches(5.3), Inches(1.2),
              items[0], sz=12, color=DARK_TEXT, lsp=1.25)
    for it in items[1:]:
        ap(tf_m, it, sz=12, color=DARK_TEXT, sa=1)

add_notes(s,
"""For teaching, I plan to inherit three existing undergraduate courses: Building Environmental Engineering, Architectural Acoustics, and Probability and Statistics.
I also propose developing three new courses. First, Smart Building Environment Systems for graduate students, integrating IoT, machine learning, and VR. Second, AI for Built Environment, a hands-on lab course using Python and MATLAB for acoustic data analysis. Third, Sustainable Building and Ecology at the undergraduate level, covering LEED and WELL standards.
For student mentoring, undergraduates will engage in capstone designs linked to real industry projects. Graduate students will target one SCI paper and one patent each, with mandatory international conference presentations and short-term visits to our partner institutions. Career paths include Hyundai Motor, Samsung Research, construction firms, and government research institutes.""",
"""교육 계획입니다. 기존 학부 교과목인 건축환경공학, 건축음향, 확률통계론을 계승합니다.
신규 교과목으로 세 과목을 개설합니다. 대학원 스마트 건축환경 시스템(IoT+ML+VR), 대학원 AI for Built Environment(Python/MATLAB 실습), 학부 지속가능 건축(LEED/WELL)입니다.
학생 지도에서, 학부생은 산학과제 연계 캡스톤 디자인을, 대학원생은 1인 1 SCI + 1 특허를 목표로 하며, 국제학회 발표를 의무화하고 파트너 기관 단기 방문을 지원합니다.
졸업 후 진로는 현대차, 삼성리서치, 건설사, 정부 연구기관 등으로 연계합니다.""")


# ═══════════════════════════════════════
# S19. CONTRIBUTION
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
content_header(s, "Contribution to the Department")
slide_num(s, 19)

contribs = [
    ("01", "Completing Building\nEnvironment Research",
     "Sound joins Thermal + Air Quality\nfor a comprehensive research group.\n→ Joint BK21 / Research Center applications."),
    ("02", "Expanding Industry\nCollaboration",
     "Hyundai Motor + Samsung networks\nopen new partnership channels.\n→ Acoustics consulting for construction firms."),
    ("03", "Elevating Global\nPresence",
     "UCL / RWTH / Sorbonne partnerships.\n→ Department-level MOU + exchanges.\n→ ISO 12913 standardization leadership."),
    ("04", "Hanyang Alumni\nCommitment",
     "B.S. + Ph.D. from this department.\nDeep understanding of culture & history.\n→ Continuing the tradition of building\n   acoustics research at Hanyang."),
]
for i, (num, title, desc) in enumerate(contribs):
    col, row = i % 2, i // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(1.2) + Inches(row * 3.0)
    add_rounded_shape(s, x, y, Inches(6.0), Inches(2.65), LIGHT_CARD)
    nb = add_rounded_shape(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), SAMSUNG_DARK)
    tf_n = nb.text_frame; tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.paragraphs[0].text = num; tf_n.paragraphs[0].font.size = Pt(15)
    tf_n.paragraphs[0].font.color.rgb = WHITE; tf_n.paragraphs[0].font.bold = True
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(s, x + Inches(0.9), y + Inches(0.15), Inches(4.8), Inches(0.6),
       title, sz=15, color=SAMSUNG_DARK, bold=True, lsp=1.1)
    tb(s, x + Inches(0.3), y + Inches(0.85), Inches(5.4), Inches(1.6),
       desc, sz=12, color=MID_GRAY, lsp=1.3)

add_notes(s,
"""I see four ways I can contribute to this department.
First, completing the building environment research group. Adding sound to the existing thermal and air quality research enables joint applications for BK21 and Leading Research Center grants.
Second, expanding industry collaboration. My networks at Hyundai Motor and Samsung open new partnership channels for the entire department, along with acoustics consulting for construction firms.
Third, elevating global presence. My active collaborations with UCL, RWTH Aachen, and Sorbonne can be leveraged for department-level MOUs and student exchanges. I can also lead ISO 12913 standardization efforts.
Fourth, as a Hanyang alumnus who completed both B.S. and Ph.D. here, I have a deep understanding of this department's culture and history. I am committed to continuing the tradition of building acoustics research that has been established at Hanyang.""",
"""학과에 대한 네 가지 기여 방안을 말씀드립니다.
첫째, 건축환경 연구 그룹 완성입니다. 기존 열/공기 연구에 소리를 더하여 BK21, 선도연구센터 등 공동 신청의 기반을 마련합니다.
둘째, 산학협력 확대입니다. 현대차, 삼성 네트워크를 통해 학과 전체에 새로운 산학 기회를 제공하고, 건설사 음향 컨설팅도 가능합니다.
셋째, 국제 위상 제고입니다. UCL/RWTH/소르본과의 파트너십을 학과 차원 MOU와 학생 교류로 확대하고, ISO 12913 표준화를 주도합니다.
넷째, 한양대 동문으로서의 헌신입니다. 이 학과에서 학사와 박사를 모두 마쳤기에, 학과 문화와 역사를 깊이 이해하고 있으며, 한양대 건축음향 연구의 전통을 이어가겠습니다.""")


# ═══════════════════════════════════════
# S20. THANK YOU
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SAMSUNG_DARK)

tb(s, Inches(1.5), Inches(1.2), Inches(10.3), Inches(2.8),
   "From Hanyang's tradition\nin building acoustics,\nthrough industry validation,\n"
   "to intelligent sound environments\n— SENSE Lab is ready.",
   sz=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER, lsp=1.4)

add_shape(s, Inches(5.5), Inches(4.3), Inches(2.3), Pt(2), SAMSUNG_CYAN)

tb(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
   "Hyun In Jo, Ph.D.", sz=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.35),
   "best2012@naver.com  |  linkedin.com/in/hyunin-jo",
   sz=13, color=SAMSUNG_CYAN, align=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
   "Department of Architectural Engineering\nHanyang University",
   sz=15, color=RGBColor(0x77, 0x88, 0x99), align=PP_ALIGN.CENTER, lsp=1.3)
tb(s, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.4),
   "THANK YOU", sz=18, color=SAMSUNG_CYAN, bold=True, align=PP_ALIGN.CENTER)

add_notes(s,
"""Thank you for your attention. From Hanyang's tradition in building acoustics, through industry validation at Hyundai Motor Company, to the next generation of intelligent sound environments — SENSE Lab is ready.
I would be happy to answer any questions you may have.""",
"""경청해 주셔서 감사합니다. 한양대학교 건축음향의 전통 위에, 산업체에서의 검증을 거쳐, 차세대 지능형 음환경 연구로 — SENSE Lab은 준비되어 있습니다.
질문이 있으시면 기꺼이 답변드리겠습니다.""")


# ─── Save ───
out = "/Users/hyunbin/Research/HanyangUniv_Faculty_Presentation_HyunInJo.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
