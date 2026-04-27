"""
pptx_helpers.py -- Samsung Research Portfolio PPT design system module.
Modern Glass design language: gradient backgrounds, rounded rectangles,
glass cards with subtle borders, pill badges, and gradient accent bars.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Samsung Brand Color Constants
# ---------------------------------------------------------------------------
DEEP_NAVY   = RGBColor(0x0B, 0x1D, 0x6F)
NAVY        = RGBColor(0x13, 0x28, 0x9F)
MID_NAVY    = RGBColor(0x0F, 0x22, 0x87)  # gradient midpoint
ACCENT_BLUE = RGBColor(0x3D, 0x7D, 0xDE)
SKY_BLUE    = RGBColor(0x06, 0x89, 0xD8)
DARK_TEXT    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY         = RGBColor(0x5B, 0x71, 0x8D)
GRAY2        = RGBColor(0x75, 0x78, 0x7B)
LIGHT_BG     = RGBColor(0xFA, 0xFB, 0xFD)
CARD_BORDER  = RGBColor(0xE7, 0xE6, 0xE2)
WARM_GRAY    = RGBColor(0xE7, 0xE6, 0xE2)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
SHADOW_COLOR = RGBColor(0xD8, 0xDA, 0xE0)
RED_TINT     = RGBColor(0xFD, 0xF0, 0xEF)
GREEN_TINT   = RGBColor(0xEF, 0xF9, 0xF2)
RED_TEXT     = RGBColor(0xC0, 0x39, 0x2B)
GREEN_TEXT   = RGBColor(0x27, 0xAE, 0x60)

# ---------------------------------------------------------------------------
# Slide Dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Margins
# ---------------------------------------------------------------------------
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.5)

# ---------------------------------------------------------------------------
# Font Constants
# ---------------------------------------------------------------------------
FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"
FONT_EN    = "Segoe UI"

# ---------------------------------------------------------------------------
# Presentation Factory
# ---------------------------------------------------------------------------

def new_presentation() -> Presentation:
    """Return a blank 16:9 Presentation with Samsung dimensions."""
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

# ---------------------------------------------------------------------------
# Slide Helpers
# ---------------------------------------------------------------------------

def add_blank_slide(prs: Presentation):
    """Add and return a fully blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_slide_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_gradient_bg(slide, color1: RGBColor, color2: RGBColor):
    """
    Set a two-stop linear gradient background on a slide.
    Uses direct XML manipulation for gradient backgrounds.
    Falls back to midpoint solid if gradient XML fails.
    """
    try:
        bg = slide.background
        bgPr = bg._element
        # Remove any existing bgPr children that are fill-related
        for child in list(bgPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('bgFill', 'bg'):
                bgPr.remove(child)

        # Build gradient fill XML
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        bgFill_el = bgPr.find(qn('p:bgPr'))
        if bgFill_el is None:
            bgFill_el = etree.SubElement(bgPr, qn('p:bgPr'))

        # Remove existing fills in bgPr
        for child in list(bgFill_el):
            bgFill_el.remove(child)

        gradFill = etree.SubElement(bgFill_el, qn('a:gradFill'))
        gradFill.set('flip', 'none')
        gradFill.set('rotWithShape', '1')

        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        # Stop 1 (position 0%)
        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '0')
        srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
        srgb1.set('val', '%02X%02X%02X' % (color1[0] if isinstance(color1, (list, tuple)) else int(str(color1)[:2], 16),
                                              color1[1] if isinstance(color1, (list, tuple)) else int(str(color1)[2:4], 16),
                                              color1[2] if isinstance(color1, (list, tuple)) else int(str(color1)[4:6], 16)))

        # Stop 2 (position 100%)
        gs2 = etree.SubElement(gsLst, qn('a:gs'))
        gs2.set('pos', '100000')
        srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
        srgb2.set('val', '%02X%02X%02X' % (color2[0] if isinstance(color2, (list, tuple)) else int(str(color2)[:2], 16),
                                              color2[1] if isinstance(color2, (list, tuple)) else int(str(color2)[2:4], 16),
                                              color2[2] if isinstance(color2, (list, tuple)) else int(str(color2)[4:6], 16)))

        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', '5400000')  # top to bottom
        lin.set('scaled', '1')

        # Ensure <a:effectLst/> is there
        etree.SubElement(bgFill_el, qn('a:effectLst'))
    except Exception:
        # Fallback: use solid midpoint color
        r = (color1[0] + color2[0]) // 2 if isinstance(color1, (list, tuple)) else 0x0F
        g = (color1[1] + color2[1]) // 2 if isinstance(color1, (list, tuple)) else 0x22
        b = (color1[2] + color2[2]) // 2 if isinstance(color1, (list, tuple)) else 0x87
        set_slide_bg(slide, RGBColor(r, g, b))


def _set_gradient_bg_simple(slide, color: RGBColor):
    """Simple solid fill for gradient slides (reliable fallback)."""
    set_slide_bg(slide, color)


# ---------------------------------------------------------------------------
# Shape Primitives
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height,
             fill_color: RGBColor = None,
             line_color: RGBColor = None,
             line_width_pt: float = 0.0):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    fill = shape.fill
    if fill_color is not None:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color is not None:
        line.color.rgb = line_color
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()

    return shape


def add_rounded_rect(slide, left, top, width, height,
                     fill_color: RGBColor = None,
                     border_color: RGBColor = None,
                     border_width_pt: float = 1.0):
    """Add a rounded rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    fill = shape.fill
    if fill_color is not None:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if border_color is not None:
        line.color.rgb = border_color
        line.width = Pt(border_width_pt)
    else:
        line.fill.background()

    # Set corner radius to a moderate value
    try:
        sp = shape._element
        sp_pr = sp.find(qn('a:prstGeom'))
        if sp_pr is not None:
            avLst = sp_pr.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(sp_pr, qn('a:avLst'))
            for gd in list(avLst):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 8000')  # moderate rounding
    except Exception:
        pass

    return shape


def add_glass_card(slide, left, top, width, height,
                   fill_color: RGBColor = None,
                   border_color: RGBColor = None,
                   shadow: bool = True):
    """
    White card with subtle border simulating glass effect.
    Adds a slight shadow rectangle behind it for depth.
    Returns the main card shape.
    """
    fill_color = fill_color or WHITE
    border_color = border_color or CARD_BORDER

    # Shadow layer (offset down-right by 2pt)
    if shadow:
        add_rounded_rect(slide,
                         left + Pt(3), top + Pt(3),
                         width, height,
                         fill_color=SHADOW_COLOR,
                         border_color=None)

    # Main card
    card = add_rounded_rect(slide, left, top, width, height,
                            fill_color=fill_color,
                            border_color=border_color,
                            border_width_pt=0.75)
    return card


def add_accent_bar(slide, left, top, width, height=Inches(0.04),
                   color: RGBColor = None):
    """Add a thin horizontal accent bar (defaults to NAVY)."""
    color = color or NAVY
    return add_rect(slide, left, top, width, height, fill_color=color)


def add_gradient_bar(slide, top, height=Pt(4)):
    """
    Full-width gradient accent bar at top of content slides.
    Simulated with 3 adjacent colored rectangles for gradient feel.
    """
    third = int(SLIDE_WIDTH) // 3
    add_rect(slide, 0, top, Emu(third), height, fill_color=DEEP_NAVY)
    add_rect(slide, Emu(third), top, Emu(third), height, fill_color=NAVY)
    add_rect(slide, Emu(third * 2), top, Emu(third + 10), height, fill_color=ACCENT_BLUE)


def add_pill_badge(slide, left, top, text,
                   bg_color: RGBColor = None,
                   text_color: RGBColor = None,
                   font_size: float = 9):
    """Small rounded pill for paper references."""
    bg_color = bg_color or NAVY
    text_color = text_color or WHITE

    # Calculate width based on text length
    char_width = font_size * 0.55
    pill_w = max(Pt(len(text) * char_width + 30), Inches(1.5))
    pill_h = Inches(0.32)

    pill = add_rounded_rect(slide, left, top, pill_w, pill_h,
                            fill_color=bg_color)
    add_textbox(slide, left, top, pill_w, pill_h,
                text, font_size=font_size, font_color=text_color,
                bold=True, alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_EN)
    return pill


# ---------------------------------------------------------------------------
# Text Helpers
# ---------------------------------------------------------------------------

def add_textbox(slide, left, top, width, height, text: str,
                font_size: float = 12,
                font_color: RGBColor = None,
                bold: bool = False,
                alignment=PP_ALIGN.LEFT,
                font_name: str = None,
                anchor=MSO_ANCHOR.TOP) -> object:
    """Add a single-paragraph textbox and return the shape."""
    font_color = font_color if font_color is not None else DARK_TEXT
    font_name  = font_name  or FONT_BODY

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = alignment

    run = p.add_run()
    run.text = text

    fnt = run.font
    fnt.size  = Pt(font_size)
    fnt.bold  = bold
    fnt.color.rgb = font_color
    fnt.name  = font_name

    return txBox


def add_multiline_textbox(slide, left, top, width, height,
                          lines,
                          font_size: float = 11,
                          font_color: RGBColor = None,
                          line_spacing: float = 1.15,
                          font_name: str = None,
                          alignment=PP_ALIGN.LEFT) -> object:
    """
    Add a textbox with multiple lines/runs.

    lines : list of items, each either:
            - str
            - tuple (text, bold, color, size)
            - tuple (text, bold, color)
            - tuple (text, bold)
    """
    font_color = font_color if font_color is not None else DARK_TEXT
    font_name  = font_name  or FONT_BODY

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    first = True
    for item in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        para.alignment = alignment

        # Parse line spec
        if isinstance(item, str):
            t, b, c, s = item, False, font_color, font_size
        elif isinstance(item, (list, tuple)):
            if len(item) == 4:
                t, b, c, s = item
            elif len(item) == 3:
                t, b, c = item; s = font_size
            elif len(item) == 2:
                t, b = item; c = font_color; s = font_size
            else:
                t = item[0]; b = False; c = font_color; s = font_size
        else:
            t, b, c, s = str(item), False, font_color, font_size

        # Line spacing
        pPr = para._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))

        run = para.add_run()
        run.text = t

        fnt = run.font
        fnt.size  = Pt(s)
        fnt.bold  = b
        fnt.color.rgb = c if c is not None else font_color
        fnt.name  = font_name

    return txBox


# ---------------------------------------------------------------------------
# Composite Components
# ---------------------------------------------------------------------------

def add_section_title(slide, left, top, width,
                      title_text: str,
                      subtitle_text: str = "",
                      title_size: float = 28,
                      subtitle_size: float = 13):
    """
    Add a section title block: gradient bar + big title + optional subtitle.
    Returns (title_shape, subtitle_shape or None).
    """
    # Gradient accent bar at the very top of the slide
    add_gradient_bar(slide, top=Inches(0), height=Pt(5))

    title_top = top + Inches(0.15)
    title_h   = Inches(0.55)
    t_shape   = add_textbox(
        slide, left, title_top, width, title_h,
        title_text,
        font_size=title_size,
        font_color=NAVY,
        bold=True,
        font_name=FONT_TITLE,
    )

    s_shape = None
    if subtitle_text:
        sub_top = title_top + title_h
        s_shape = add_textbox(
            slide, left, sub_top, width, Inches(0.35),
            subtitle_text,
            font_size=subtitle_size,
            font_color=GRAY,
            bold=False,
            font_name=FONT_BODY,
        )

    return t_shape, s_shape


def add_metric_card(slide, left, top, width, height,
                    number_text: str,
                    label_text: str,
                    bg_color: RGBColor = None,
                    number_size: float = 28,
                    label_size: float = 10):
    """
    A glass-style metric card: white bg, subtle border, shadow,
    big navy number on top, small gray label below.
    """
    bg_color = bg_color or WHITE

    # Glass card with shadow
    card = add_glass_card(slide, left, top, width, height,
                          fill_color=bg_color, shadow=True)

    pad     = Inches(0.12)
    inner_w = width - pad * 2

    # Number
    num_h   = int(height * 0.55)
    add_textbox(
        slide,
        left + pad,
        top  + pad,
        inner_w,
        Emu(num_h),
        number_text,
        font_size=number_size,
        font_color=NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name=FONT_EN,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Label
    lbl_top = top + Emu(num_h)
    lbl_h   = height - Emu(num_h) - pad
    add_textbox(
        slide,
        left + pad,
        lbl_top,
        inner_w,
        lbl_h,
        label_text,
        font_size=label_size,
        font_color=GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
        font_name=FONT_BODY,
        anchor=MSO_ANCHOR.TOP,
    )

    return card


def add_implication_box(slide, left, top, width, height,
                        lines,
                        bg_color: RGBColor = None,
                        font_size: float = 11,
                        line_spacing: float = 1.3):
    """
    Navy rounded background box with white text lines.
    Returns the background shape.
    """
    bg_color = bg_color or NAVY

    bg = add_rounded_rect(slide, left, top, width, height,
                          fill_color=bg_color)

    pad    = Inches(0.18)
    inner_w = width  - pad * 2
    inner_h = height - pad * 2

    # Normalise lines so all colors default to WHITE
    normalised = []
    for item in lines:
        if isinstance(item, str):
            normalised.append((item, False, WHITE, font_size))
        elif isinstance(item, (list, tuple)):
            if len(item) == 4:
                normalised.append(item)
            elif len(item) == 3:
                normalised.append((item[0], item[1], item[2], font_size))
            elif len(item) == 2:
                normalised.append((item[0], item[1], WHITE, font_size))
            else:
                normalised.append((str(item[0]), False, WHITE, font_size))
        else:
            normalised.append((str(item), False, WHITE, font_size))

    add_multiline_textbox(
        slide,
        left  + pad,
        top   + pad,
        inner_w,
        inner_h,
        normalised,
        font_size=font_size,
        font_color=WHITE,
        line_spacing=line_spacing,
        font_name=FONT_BODY,
        alignment=PP_ALIGN.LEFT,
    )

    return bg


def add_image_safe(slide, image_path: str, left, top, width, height,
                   placeholder_color: RGBColor = None,
                   placeholder_text: str = ""):
    """
    Add an image to the slide. If the file is missing, add a glass-card
    placeholder with label text instead.
    """
    if image_path and os.path.isfile(image_path):
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        return pic
    else:
        ph_color = placeholder_color or LIGHT_BG
        rect = add_glass_card(slide, left, top, width, height,
                              fill_color=ph_color, shadow=False)
        label = placeholder_text or (os.path.basename(image_path) if image_path else "Image")
        add_textbox(
            slide, left, top, width, height,
            label,
            font_size=10,
            font_color=GRAY,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        return rect
