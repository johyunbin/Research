"""
extract_images.py
Extracts figures from PPT/PPTX and PDF files into /Users/hyunbin/Research/assets/
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pdf2image import convert_from_path

ASSETS_DIR = Path("/Users/hyunbin/Research/assets")
POPPLER_PATH = "/opt/homebrew/bin"


def extract_pptx_images(pptx_path: str, slide_indices: list, prefix: str) -> list:
    """
    Extract all PICTURE shapes from specified slides in a PPTX file.

    Args:
        pptx_path: Path to the .pptx file
        slide_indices: 1-based slide numbers to extract from
        prefix: Filename prefix for saved images

    Returns:
        List of saved file paths
    """
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    saved = []

    for slide_num in slide_indices:
        idx = slide_num - 1  # convert to 0-based
        if idx < 0 or idx >= len(prs.slides):
            print(f"  [WARN] Slide {slide_num} out of range (total: {len(prs.slides)}), skipping")
            continue

        slide = prs.slides[idx]
        img_count = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_count += 1
                image = shape.image
                ext = image.ext  # e.g. 'png', 'jpeg'
                filename = f"{prefix}_s{slide_num}_img{img_count}.{ext}"
                out_path = ASSETS_DIR / filename
                with open(out_path, "wb") as f:
                    f.write(image.blob)
                print(f"  Saved: {out_path}")
                saved.append(str(out_path))

        if img_count == 0:
            print(f"  [INFO] Slide {slide_num}: no PICTURE shapes found")

    return saved


def extract_pdf_pages(pdf_path: str, pages: list, prefix: str, dpi: int = 200) -> list:
    """
    Render specific PDF pages as PNG images.

    Args:
        pdf_path: Path to the .pdf file
        pages: 1-based page numbers to render
        prefix: Filename prefix for saved images
        dpi: Resolution for rendering (default 200)

    Returns:
        List of saved file paths
    """
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    saved = []

    for page_num in pages:
        images = convert_from_path(
            pdf_path,
            dpi=dpi,
            first_page=page_num,
            last_page=page_num,
            poppler_path=POPPLER_PATH,
        )
        if not images:
            print(f"  [WARN] Page {page_num} could not be rendered from {pdf_path}")
            continue

        filename = f"{prefix}_p{page_num}.png"
        out_path = ASSETS_DIR / filename
        images[0].save(str(out_path), "PNG")
        print(f"  Saved: {out_path}")
        saved.append(str(out_path))

    return saved


def main():
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    print(f"Assets directory: {ASSETS_DIR}\n")

    # ------------------------------------------------------------------ #
    # 1. Portfolio PPT  (Portfolio_HyunInJo.pptx)
    # ------------------------------------------------------------------ #
    portfolio_pptx = "/Users/hyunbin/Research/Portfolio_HyunInJo.pptx"
    print("=== Portfolio PPT ===")

    print("  Slides 5 (Pleasant-Eventful diagram), prefix=proto")
    extract_pptx_images(portfolio_pptx, [5], "proto")

    print("  Slides 10, 11, 18 (HRTF results), prefix=proto")
    extract_pptx_images(portfolio_pptx, [10, 11, 18], "proto")

    print("  Slides 24, 25 (AI/LSTM figures), prefix=proto")
    extract_pptx_images(portfolio_pptx, [24, 25], "proto")

    # ------------------------------------------------------------------ #
    # 2. HMG PPT
    # ------------------------------------------------------------------ #
    hmg_pptx = (
        "/Users/hyunbin/Research/ETC/"
        "HMG 학술대회 발표자료_AVAS_사운드스케이프_MSV소음진동시험팀_조현인책임_충남대공유용.pptx"
    )
    print("\n=== HMG PPT ===")
    print("  Slides 10, 11, 12, 13, 17 (AVAS results), prefix=hmg")
    extract_pptx_images(hmg_pptx, [10, 11, 12, 13, 17], "hmg")

    # ------------------------------------------------------------------ #
    # 3. Research paper PDFs
    # ------------------------------------------------------------------ #
    paper_dir = Path("/Users/hyunbin/Research/Paper")

    pdf_tasks = [
        (
            "APAC_2022_Jo&Jeon_Perception of urban soundscape and landscape using different visual environment reproduction methods in virtual reality.pdf",
            list(range(5, 9)),   # pages 5-8
            "apac2022",
        ),
        (
            "B&E_2019_Jeon&Jo_Three-dimensional virtual reality-based subjective evaluation of road traffic noise in urban high-rise residential buildings.pdf",
            list(range(5, 11)),  # pages 5-10
            "be2019a",
        ),
        (
            "B&E_2020_Jeon&Jo_Effects of audio-visual interactions on soundscape and landscape perception and their influence on satisfaction with the ur.pdf",
            list(range(6, 11)),  # pages 6-10
            "be2020",
        ),
        (
            "SCS_2021_Jo&Jeon_Compatibility of quantitative and qualitative data-collection protocols for urban soundscape evaluation.pdf",
            [4, 5, 15, 16, 17],  # pages 4-5, 15-17
            "scs2021",
        ),
        (
            "SCS_2023_Jeon et al_Psycho-physiological restoration with audio-visual interactions through virtual reality simulations of soundscape and land.pdf",
            list(range(7, 12)),  # pages 7-11
            "scs2023",
        ),
        (
            "IJERPH_2024_Jo&Jeon_Quantification of visual attention by using eye-tracking technology for soundscape assessment through physiological response.pdf",
            list(range(4, 8)),   # pages 4-7
            "ijerph2024",
        ),
        (
            "B&E_2021_Jo&Jeon_Overall environmental assessment in urban park_Modelling audio-visual interaction with a structural equation model based on.pdf",
            list(range(8, 12)),  # pages 8-11
            "be2021",
        ),
        (
            "SENSORS_2021_Chung et al_Diagnosis of pneumonia by cough sounds analyzed with statistical features and AI.pdf",
            list(range(3, 8)),   # pages 3-7
            "sensors2021",
        ),
    ]

    print("\n=== Research Paper PDFs ===")
    for filename, pages, prefix in pdf_tasks:
        pdf_path = paper_dir / filename
        if not pdf_path.exists():
            print(f"  [WARN] File not found: {pdf_path}")
            continue
        print(f"  [{prefix}] {filename}  pages={pages}")
        extract_pdf_pages(str(pdf_path), pages, prefix)

    # ------------------------------------------------------------------ #
    # Summary
    # ------------------------------------------------------------------ #
    all_assets = list(ASSETS_DIR.iterdir())
    print(f"\nDone. Total assets extracted: {len(all_assets)}")
    for f in sorted(all_assets):
        size_kb = f.stat().st_size // 1024
        print(f"  {f.name}  ({size_kb} KB)")


if __name__ == "__main__":
    main()
